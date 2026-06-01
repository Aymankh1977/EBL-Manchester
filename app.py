"""
DentEdTech™ Evidence Engine
An educational AI platform for medicine, dentistry, and pharmacology students
at Manchester University. Built on the REAL-AI framework principles.

© DentEdTech™ - All Rights Reserved
"""

import streamlit as st
import anthropic
import json
import re
import smtplib
import uuid
import base64
from io import BytesIO
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime


# ─── Embedded Logo ───
LOGO_B64 = "/9j/4AAQSkZJRgABAQAAAQABAAD/2wBDAAgGBgcGBQgHBwcJCQgKDBQNDAsLDBkSEw8UHRofHh0aHBwgJC4nICIsIxwcKDcpLDAxNDQ0Hyc5PTgyPC4zNDL/2wBDAQkJCQwLDBgNDRgyIRwhMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjIyMjL/wgARCAKAAoADASIAAhEBAxEB/8QAGgABAAMBAQEAAAAAAAAAAAAAAAEDBAUCBv/EABgBAQEBAQEAAAAAAAAAAAAAAAABAgME/9oADAMBAAIQAxAAAALkhAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAB7rx73Xb5c+OkueRHYozvnxpz53CYlATEkT66BzbevSY41XnGj6LmHPAAAAAAAAAAAAAAAAAAAAAAAAAAAB66ObX14zPmdc5QWUElBaqtVedUZvfrn1p36sk1s5+WAbTFPX5hPX4Xs28/vcQ8AAAAAAAAAAAAAAAAAAAAAAAAAAE10PZ180oVKBKB5p0M6xeN3uarupwZ37qTnUJgfQfP8ASKNezimcG6cvQOUAAAAAAAAAAAAAAAAAAAAAAAAAB78e7Ogh288oEomJQqUCaGPn1QY6XdXFoLeN2eSVevPo6nJ6/IC28ydfj9g5EAAAAAAAAAAAAAAAAAAAAAAAAAA9eVdKInr55QJQJQJ8zBgjRXy7V6M/Zmoy4/Jqyg6QRzg3beLJPTwazmgAAAAAAAAAAAAAAAAAAAAAAAAAA2+89/XhKGpKBKJiWCMdOhmoSx1uT1c75cA6UeD1zwm/VBg89XlG7xr5Z5AAAAAAAAAAAAAAAAAAAAAAAAAAB714de+fsb5pgSgZY1+M9Mg59HS5voto6eE6HK6fMEumWU56zbhr6B65l1IAAAAAAAAAAAAAAAAAAAAAAAAAAAtqmzYh14ygSgSgK7EZI1Z8dXQ5vrOujzOlhLbPAygu1sNeBAAAAAAAAAAAAAAAAAAAAAAAAAAAAGj3Vb14hYAAnzRnV8ZmdTHrzN3X0XHgGY9R0Ob0OeAAAAAAAAAAAAAAAAAAAAAAAAAAAAAW20XdOUoakoEo8S+PBy6gAW3U3V4NRzbZ9Rbi05gAAAAAAAAAAAAAAAAAAAAAAAAAAAACb89+8SN4lAmi2jGwxsAC22qyq7siLL8m+stUxAAAAAAAAAAAAAAAAAAAAAAAAAAAAACyv1ZaOnMDzV78c+gSgAW2V+7M4l9aq/NlIlAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAtmuzpzCyvzMcuoAAFvujVZkmNK+8llQEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAPfhZa8NZ8jGwAAGnN6q+K60CUAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA0a9Y5jpxXNb80tKYzoAAAAAezwvoobExuo1nluoOW1Zc6CU9da547R1prgu8OC7w4LdhAAAAC/unza2oAAAAAAAAAAAAA0acujpzyTUx01aObs1zpp2Y86TpvsweejWYk6ZrLO71rOHx0qjFor3RPM6TWebr8es6vyaMJcpTXvwsPHrb61nn9Ov1Zi9eKsdNHV4XajLROUs8X9M5HjswcZbWR66Gw4vju5jB2+J2Tix1rTgx3eQUgAAAAAAAAAAAu0Z7+nPEOfRpzadZszX13OnDt82Y9sSuf3CW7HpyS268Wy5xW1WTWrBtwpOjNoXQrzaztYksa8W88ZLac6s28/frOOuyvO3Z43YMOXTSdXHt8nI6NsxVm10G7mbuST7rFvY4/WOfl91nZy21GAAAAAAAAAAAAFt1N3TlkaZms+qK7I81zjevJobxmX+s3NootWzLpzHrXk1XOX34nO9WLX43jPf49ZtmXXGplakuXZQj1Rr8WUba/Vmev34x1dbk9Qx0X5zr8+zUcpumMGrxVXV5XSmOU6EVl6nK6kcuv34OnVZVWIQAAAAAAAAAAB6nws9vA9eSUB78KsjwSZ8l9+CJ91qmCPc1rPTyWyPBPbwPXkl9e6lnuayzBC6ke/ALahpjOPXkJvzq011I9XZxMBd5rAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAH//aAAwDAQACAAMAAAAh/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wDPIdlt/wCu9N2H/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP3KRDbHn1d/72X/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A4zEc8qdQ73vD/n//AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AKTfewQWnja73923/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wCxoIL7lvu/+99cv/8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/Ps87DD/AH+323x//wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AKzjTE/7bf8A4/3/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/rGOOOf5P7j/ABf/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD+s88439JMO/8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A88wwx7//AIwjf/8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP6xBd//AP5HmN//AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD+tsn/AP8A/r/5/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/ANHP7/8A/J4s/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP3n/wD/AP8Ah/8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD89PN//wD/AP8A+tOvPNf84t89/wD/AP8A5/8A/wD/AP8A/wD/AP8A/wD/AP8A8Tz9/fvaxy938d3kxhSees/POvc/P/8A/wD/AP8A/wD/AP8A/wD4P7M1VNPkqR3NDCBuPOEaLh5z516//wD/AP8A/wD/AP8A/wD/APx4k6szuU7mc3Qwjf8AXPJUN+M+Wn4P/wD/AP8A/wD/AP8A/wD/APd9f/8APXHjHnXHfX7TP/n7j37Lj/3/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/wD/AP8A/9oADAMBAAIAAwAAABD77777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777777j1PZ3b6qxUsf7777777777777777777777777775ivrMJ3Oa3igrX777777777777777777777777771ouc+4bLrY2X63777777777777777777777777767EXo01F8/HpX7l77777777777777777777777776hIJ4IP+DT6VyED777777777777777777777777775JUV4Tsx6n6FRX777777777777777777777777775/KYTz6qJaK2Z7777777777777777777777777776zDLLLdr56r6z77777777777777777777777777776vDDCr9VUnr777777777777777777777777777775PLLJ/77nnbb77777777777777777777777777775nwwX777x+B777777777777777777777777777776vv5777436/7777777777777777777777777777777f1b7767jz7777777777777777777777777777777+f77775z7777777777777777777777777777777777777777777777777777777777777778bfb7777776kIf36PMoxb7777D77777777777774iyx/aa1kGKDlmN0gBz45KrrKbJrL777777777776hbojMhwRPl+8z5lU7QFyU/647q5r777777777776M+xtX+TeJoz/wBxD9HUF8pypis/+N++++++++++++dve+8+N/s/+sf/APr7Hn3v7nvfzvjvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvvv/EACwRAAEEAAQGAQQCAwAAAAAAAAEAAgMREBIhMQQUIEBBUTIiM0JQYYATMGD/2gAIAQIBAT8A/s61pcaCZwY3cUeFjTuDr4lOic3cdFq1ZQPe8NGGts7nC8N1I2Jo+oKQtJ+gI2gESgUQh3bRZpAUKVq1aexx2NLl9bcbUpaNG4lDDz3cfyCtWrV4Tz/i3DyhgMD3bTTr6X3lNJ8TmiyEVWFoKu9idmaDiVzZ9KWfOKpHpHe8K/dp6H8OwNJ6q76N2VwKHQ5gcKIUsRZr4/Rxm2g9Es4Zp5R4gn5DRGr06D30DrYFatSPyNJKJLjZx8KtF5R77hnbhWrXEu1A6PCspvfwuyvBxlNuPR+OGzf0ET87cHGyT0DVpCa2ynGz+gY8sNhOnbl036WuopzhVD/gGQ5m5rXLD2nwluvWASaCc0t3CjbndS5Ye1y49pwpxGDI8wu0RRrC+q/98f2lmPtQSFxoqVtP0TYBVuNJ0Aq2lRszupcuL1OifAALaVBH+SkjDjuohUlKZxzaFZne0AXGkIG+SmMy2E7c9uz7WHDj6k4AyhSR5iBaiYGXqo6EppTk5qXDn6CovuBcRuFB8k+bIapcwPSgF2VKSXEKEktKdue3iFx0uW/lDJEN0JDmzKRglAIKHDuG5XD/ADXEfNQfAqM08EqWEvoqEZZKUkOd1grlv5UZyOLU+DMcwUbMgTtyhjWFdiHuqgVnd7RN74BxbsUZHHcoEg2ESTqSg4jY4f5HAUCsxu71Wd3tZ3e1qTZWd2wKzne1d9F/2d//xAAxEQABBAAEAwYFBAMAAAAAAAABAAIDEQQQEiEgMUATFDIzQVEiUGFxkQVCgKFggbH/2gAIAQMBAT8A/k65wYNTjspP1Ag1GPyhj5r3pR/qAOzxSZNG/wAJvgpUtlXW42YvfpHIf9VZUh9FE+dxphJ/tRNe0fGbKCJypAo9W4020TZJKrKlHIxvibf+13yhpY2lCHH4nZhHL06uXwO+ypUqVZYXDV8b8hyRyKpDq3i2kKq4I6DgXclHiI3mm8V9bMzQ9wzG67iPdQ4URODr4ayPWY2Pk8cEWKkc4A0r+Rys1sIVKlSpMkcw2CoZxIK9fkcraeQqVKlDhzJudghhWtNtO6F1vwDrsQ2nlUqUTNbgEAAKGfqid0eSHLrsW3cFUqWFbsTwHmiE7YUh10zdTCqVKEUwcH7shufkErNLsmigBwHZwKcdk0UPkD2Bw3TYHat+XC5thNabs/4BJiSx2ml3s+yjxLXmjtxucGiymva/kVLJ2Yul3w+y72fb+1G7U0HKWXQQKQNgHp5fP/C0N9liYgwamqB5dHZT8UbpgtNxRunilLJ2bdSOLJGw3TMUbpwWJlI+CuaimMYIAU7tUNrDNaWWR6rSxOc1jb9F3p17BSSF5BpM8I6eXz/xliz8CZbYDShl7MEgWppO0A2UlmAWsKBoJ9ViQBIFP5RWE8JWK8tRQaxqBXdD7rFE7BQtDWivVYqg4Um+Efbp5jUtrvn0R7Sdw2RiBZoUchhJa4I4ptbBYry1hfLH3WJ8YUrSYyB7KCbs7BCxDtcVhRYjQ2iF3z6KVvasDgFFidIohSyayDSZ4R9kUMrytX0BY0myFoZ7BAAbDIsa7mEImDkEQHCiEGgCgEWNO5GXZMJshaW1VbLQz2C0M9ggABQRjaTZC0NqqQFbDgr+Tv8A/8QAQhAAAQIDAwcKBQMCBQUBAAAAAQIDAAQREBIxEyEyM0FRcQUgIjBQUmFygZEUIzRCoUBisSRDU2CCkqAVcHOi0fD/2gAIAQEAAT8C/wCS6lpxeigx8I93R7wZZ4fZ7QQU6QI49amVfXosq/iP+nTPcH+6DyfMj+3XgqFJUg0UkpPj2222pxVEw3Lob8Tv5mIzwuVaVsunwhcm4NHpQQUmihQ+PPbaW6q6hJUYZ5LGLyv9KYvSsp3EHwxhXKrQ0ULV+IHKya52Vf7oRyhLr+4pP7oUlDqOkApJia5OuC+znHd3dsgFSgBiYbbDaLo6ggKzEVEKlWjsu8IMluc9xHwR/wAQe0LCQaJVe8YAqaAVMMcnV6T2b9ohb7Eoi7h+1MPcoPO5k9BO4c1mYcYPQV6bDEtNImU5syhimOUJW4csgdE6Q3dsSiMV+nVuPIaxx3CHX1O45k7oYlVveCe8YaZalkEj1UYmOUCeizmHe5kpIGYF9Zuo/JhfJTNOiVpO+H2Fy7lxfvvsbWppwLRiIbWial6/arMRDrZadUg4jtdlN1pI6ha7n2KVwgzlP7fuYVNOK23R4QhpbmiPUw1KIRnV0j+IdfSymqvQQ8+t89LDdzNkG8jk/wCTiG+jSJFx0zibqia6XCOVgMm0dtTbya7dcU3sVn9Y5TR8xDm8U7WxPVGhxjJo7ifax6ZDQoM6v4hSitV5RqedI8oBtIZdw+1UOPMsIKz+BjE3MmadvUokaItl1XJhtX7o5RTWWr3VdrN6xPHrH37nRTpfxbLs5d9KNm3hC1syaQNEbhBQzOs1x/dtEKSULUk4jNYnTTxEcofSnzDmYRO55Rz07Wb1qePVvPXBROl/HM5MSrLldOjSlYn5Rb6krbzkClIkpcyzJvnOTU+EPryj7ixgVWI1ieIjlH6U+YWMsqfVQYDEwuQUE1Su8d1LJz6Vz/8Abe1k5lDj1S1XW1GK1NstLKmFbkDEw8+iUaCUjP8AamG+U3UjppSv8Q/POvpu5kp3DbbJSWDro8qY5RmEqAZTnoaqNnJ6h00bTngm4LysBCBedSN6hE99KrxI7XrUdSsXkFO+Phf3/iHW8nTPWtiTkZIEDRRWFrUtRUo1J5knJ0o46M+xMTk7SrbRz7VcwrUrSUTxMSib00jwzxyir5aE7z2u0atJ6uaxTY59Af8Ax8yUlLtHHB0tg3RNzmLTR4qsxNBjHwUxSuT/ADCklCrqhQ7jZyenpLX6RPrq+E90dry6sU9QcTzXPoT5LZOWugOr0tg3RNTd6rbRzbVb7eTmhcyu05hCuUqOZm6o37YnEJdl7+1OcHwslkZKXFfMYWrKOKXvPa7arrgPUFpddGC2sCpHMlXUus5NWIFKbxEwwWVfsOBjZDp/oj5bQCogAVJiWbLLIQTnhUg2XK3iE92Jx0Ns5MYqzU8IYbyrwTsxMTjlxim1WbtltV5A6hzVK5iVFKgpJziGnUzLZSoZ9oh9ksqps2GHvoleUWAEkAYmJdhMum8rS2ndEzNZU3U6H8x8S9SmVVBNTU4xJtXG7x0lfxEy7lXiftGYdssqoqm/qbqe6IW0Do5jBFDQ2JUUKCk4whaJpog+o3Q/mlVDwskh8/gInnFVDeylbZdrKu59EYxNu3G7o0lfx21hnitR1SkhYhSSk0NiFqQq8nGELTMNEe43Q62Wl3T7xI64+WJ3XDy2yybjAO/PDrmVcKu22z8sdVfT3hBAcGMEUNLJY/1CYndBHGJLXK8sTuuHltX0ZdXgjtxrRPUE0FYUsq4Wk1xsl/qExOaCeMSetPlic1w8tiRVYHjEyfkK7ca29QtV4+HOl9emJzRRxiU1p4RN60cLGBV9MTZ+UBvPbjelz3Dmpv58vrxE3ooiU1h4QpCFHpJBiYSEu0SKCkSo+YTuETZ6SR24MxHPUaq58vrxE1gmGXA2okx8Unuqh1eUXeESo+WTvMPmrx8O3UmoHNJoD1DGuETWCeY30Gh4CCamvbrZ5q8OoY1sTOCbW03nAIfVRo+PbyTQ81eI6hjWekTP22y6cVQ+qq6bu30mo5itLqGD8z0iY+02AVIEGjbfCNvb6TQ8w4nqAbpBhwXmzYyi6LxxMPLvGmwf5BSbNnVNKqjhAaF+uzdDrlOiMf8AId4wTXqkqKTmjLKp/wB86xWKxXxivWVH62V1iuEekZt0UG4QphtX204Q6wUZxnHXIbU4aJj4Vyn2+9soBcVm2xRPdHtFE90e0UT3RFE90ROYozcxOmOMVh/6w+YRdT3E+0XUdxPtF1PcT7RdR3E+0XUdxPtHKIAcboNnVSn1jPmisTH1LvnP6yV1iuEKNEK4Rl3e/GXc70NzNTRebxsfbuKzYG0AnAEwUKGKVe3NGfCMmvuK9rZTWnhB0TwtlNBXGHVFLSiMY+Jd734j4l3vfiPiXe9+IW4pyl44W5NfcV7QNMcYrD5/qVGPjZjv/gR8dMU0/wD1EA5hE3MutPBKFUFN0fGzHf8AwIdeW8RfNaQAScwqfCMi7/hL/wBvNShS9FClcBGQeH9pftG2JX6trzRWJj6l3zmAy6RmaX7QQUmigQfH9VLaZ4QvVq4cxhV5vhmh8VaPhY0wKXl47orQQFA4GHGQvwVGBiV0VcYWgOChhICR0RmseaCxX7oZ1yeZNaSeESuirjD+oVzWm8orwhKQgdEWKSFYiKw/r1WHCBgInvqB5bJaXyxqdAfmEJSgUQKDwjKJrS+mvGHWkOiix67oeaLLl04bDvsl5NKQFOiqt262f+mwz1ES31TXmisIZQhZUB0lGtbFpS4m6sVETDOQdKdmz9TL6Z4QvQVw5ktgqHNWrhDKbzg8M8Vhays1MVoc0IVeQDEwOmDviW0Txha7iCqFKUs9IwyopcA2GxWZxXGJc/MPCDgYqd5sltFXHmq01cYaTdbEPO3BmxME1OeG3FJIz5t1j+uVYcIGAie148tjaMm2lA2ROPkryQOYY+Nkk6VoKCdGJ1N5i9tSYkkXn67E54rSHppbpzG6jcLC64W7hWSncYl/qW/NZMzThcUhJupGbNtgEg1Bz74ZcyjKFHEiOUf7Z4/qWNM8IVoHhzGNCHT8sxL6SrLqe6PaLqe6PaMImPtiX0TD2qNjesTxsXrFcYl9YeEHA2y+iqHCQ2SMYyrnfMZVzvmMqvvm2Y1npYNIWPa1VmyBgIndf/phvWo8wsuJJzoT7Rk0f4afaAkJwSBwETP0y4kP7npD5/p3OHMl/qW/NY9r3PMbJT6VHrHKGi36/qWNIwrQPC1KCr/7AzCkPKwTDJovjBzgiCKGnMYwMO6s2I1ieNi9NXGGNM8IOBtYwMO6s81CqoBh8YKsGkLHdabNkDAROa70gGhB3QDUVGBiaRdeJ2Kz8ySVRxSd4hQvJKd8KSUKKVY2sfUN8bHte55jZK/TI9Yn8G/X9SzpGy6ncIoNwsU7TDG1C73GFJCsYyKd5gISNkHMTDOBhzVmxGmLF6auMNm6uwtpJwhxITSkM4GMYup7oi6nuiLqe6IdAFKCGl3cxwsySN0JSlOAsd1irNkDARN670sln6DJq9DCkhQooVEGUb3qhMs0nZXjE0KPekJUUqChiIbdS6mo9t0KQlY6QrHwrPd/MTCEpllXUgYQxr2+NjuuX5jZK/TJiewR6/qQopwjKKjKqjKKgqJ280OqjKndBcVYFFOEFaiKGzAxlVQc5rYHFCMsd0KUVQlRThGVVGVVGVVGVVBUVY2BakxljuguKVGVVBNTW34lzw9oWsuKqq1D7iBmObcY+MV3EwZtz9ohSlLNVGpsBKTUGhgTbgxoY+MV3Ewt9xwUJzbhCTdUFDER8W7+32gm8ok4mxEw4hISKU4Q48t2l7Z/ynf/xAAtEAABAgMGBgIDAQEBAAAAAAABABEhMWEQQVFxkaEggbHB8PEwUEDR4WCgcP/aAAgBAQABPyH/AKXY4SMbl6Cr1nMCjjFaG+SiFuSMSO5AtzmCFuzIFU3gDfdsKzNwUfau7J06dFmABFQoqDqeyiJAdCnwKRuPIzAlmblAEk8jVENsgD/vqiXNIJIkKB+kZANKNvJN7m8OSIR8cSUTkx+5GG8AIW5hxPG6AMsAxUnLMQb5LUSMDH+lAYDAyIIgkgBFGYUhjR5n9KEoZKdz/qeDV2PM8LvMBefQTp86cUQoGJAvY/cRj3Qd7X4nshhvz38UCyFRMZZ0xTGQYR8fKJ0dvjnywRJJcuSa2v23bDRoiG4sF9kQtYgJBjYUZlOiIJwWIUxdb+/b5HOc52vY9pegUOAEDlKFgKHuj8VuQt+cdCc5G5mU/SYZDIcBmUYkkCg7IRue89TqVkNyYWmKyGZP4mBBfMx/PtgGBiWUofA6CCAOaAy40iBUBGD3IkIJebACSwBJopQtcNZAAdqFM7gExnknYkmCP3bRYNDBMN4G8Ptoh09eJ+IPk/0iSS5d7DPJBRLC8ptyQXJr7TOEXAyKp4aSzs89igbwp8ALgcIrcm4+22P48eqTiXtbw2mVUExE4DLXu4QJHlgGQZ4RGT2eGxQ+VjYcsiYXIzCBeZop80Tcro+2J3gHXgfiMJZwESIS5JtfJf1gIHWD+h8imJUJtkURufVG3FifdP6R0RQB8LApEIKoJYEYlMd7iNqC3ft9uGBx+GLTI7AhTAjYBG0BlWRG00eB/iswlU1WScPoP3bKIQtmfEKaWI6AmMcnQf37fLgb4I1UaqOBWxKMlsHQcHnLpmqfyER2H7sAIAElIAIE2hF+qJCRzBZH9wA9SmIudTH7eFzuMTRGcZpziU5xOqzdGRRPlegtCIAkHwmOa5YNsFLQCIB/ICEcCAWiiqjkWY7BYKTw7vRHIXx+3wgkbH4HTpyiniEREgCvBAdHiRDbn0BRmRg28va05aAACDKIEmF1EZCcLlndGrDWQmETRT4P6V/3L4MxAp7X4NpwEwCQUDhB/cIkKJ87+ovCwsA2TAARSQQ0AwCM8kBuxQGwx5eiEISUySszHK4mA7D9zEUuqx+IsQxYr0KirNhEaGNh6CAQGw+NNiEFGEbiwRMkTW2TDFoytrAMR+qycMr33QJIBcgACL/ibBncbQgGRgReOHDPRSbi8YgvCqFt3U2jOTj+ZIhGRkMB929Qh8Dp0x+xNVhNxBRDFMWGDBMXB0QAkqC8CoWxdTYZKHUNm+8OFX4BvEZjyWlETNfZ1vQrf+i86oW1dTZWoBumq7Df7w4jwvbk6XF13RFrF5dQtp6mzKJfZN/eBJs4422uN1vRbw9l49QnURqE1CwgE8mZLBPmn3huq8byeXHsT0W/PZAaJcNBeYENkIDARTCZp1964uF5fBtSt2e3AABOIepRDFeX+9gEcJy4/BLyK63tbmXHJMTJ98wnhOB8HXIpc7Wncgodl1/4MY/gMNxIIYZgsONmUUFIIIkkiZn/AAYR/BU1KbgzFnJAoFCf9P8AA3JsJZ3xNIvgQfNBiZmZnT/BgGBRfiOA0SsgMv8A1DROMQnGIT1H4jimqcYhOMQnGITjEJ3k3DCxmI1TMRqmYjVMwapmI1Tgyb46g/NbEblqmGDRNgInmbkriljAv69j5oTnAeKISJkj4a0x0DWKWyl6lepQANAIGQ4Nr6qZReBJesr1Fesr1FeopoURyDX/ABbAuRecx/M8GqMQJgiqjZAeLMBBYA65KwD9FSluz4FFgubhAkwEmkUQhzqLS8OIR67pb5dFDbCUKqg0Ki0Kg0KDI3INYzlg5OSgvAzqEZcQ3VsNAXEHYKk8qIgKBDhDghAGJkOKpFEgZAwgAmEbAB1DfdIggsXBwPBuUBRjDq0QzC4OBC2ix4DFO4kZ1ThQN1/K8aqLV9OAoXmjWywQAOUsGaZcAZqQJyKFuGxMc0QWDAheJRCLzAvBMEAoTlEYAAZHFb3tY6jXRYfwV5lEWz1CcJxROKWFacgJlNQAsgQi0RY6XpZMRaQW0dTYePidVgmXcIECsUDXTD5wvZCj6JNgFqhpJ75fsUDcFGqEFxBjRvvW0WBbSYRiMcMFGqKQT3FEipnOn5VFrOnAWsFvyiKSCABJkIp9mQeSBOExCr+E2Ld3XjUUIHZOgpQyEyCE6hAuJGiElD0j0XukSTMk815lE6hgFDAKGA0W8JrXmJRBj9ZEeIk1KBwNxEUk66LpZMR6QWyd1QIF0HtGDLGzKFEU0SxjRD9AGCCM5Oq5EAJJYCJRuC3BOqrFRVI7jrbk6MsJMSOYpxgBcMUbHQzkUAbLdPyoPVcAsZxKbrwW0CMXBtEIMDAAKXmW97ItvrZsydbxYBaRtPUR1CE9gvYL3iJJcm9Aooclm8CddF0sMyLSCkZO6EER8CnRdCCb2r0lX1NBFojqFfeJowax8GxJ15jGwvBiVuOz8nZI9VabhJAAwkE8QXRKhuBkD5eGRHExwHro9vrZtCdbpYkQKW79F514ciE6xdA2bwJ15crDMi0gpOXuiCLzoI5A4RG24d+B9KXJAPGAEIfDBbtqdeAxsPwYlFrdn5OyU4FeiQHLSToULTsRclyg4LhBHcFI3NewR1xqQMsCj10fhWzcp1uCGAmRgnTkOgpoAy3yLAxYr0C9QvQIw0E7k6vlhIXcMlIROvFlYZkWkpWXvYJxbyZMe4BCKMBc1FCechwgACN3JQGSOFfnvQ1h5woju1oORgoZrZU685jYW71K3fZ+SYdmiytFUGirjRTMjz4JRCAmxzCpUVgMlOJQ5mRooKGNgJARcsnRESFfYJZwRVUaKuWhRQRkVk6KqNFVGiydFdCFLIYDDAqjQhpCiydEY5M5tu3VDk7NC1lO4B17QoiQ5Lp0mpY+cC8FSxnhuiw9QpoMZRFbDjhV0nYTjmyCzYoABuBg3/U7/8QALRABAAIBAgQFBAMBAAMAAAAAAQARITFBUWFxkRCBobHwUMHR8SAw4UBgcKD/2gAIAQEAAT8Q/wDpcUC2Ai3w13NEHLr5P+kzSL4gN+k5GhX3f2aoMrobvSHHXSu7omvfJ9FmjO2a7KMW8eeuzr9bH33VcBxX5cKHy7dPJoe/iYS+a9S47Ms1Leg7rPQlzs6HsrXZnI4iL1/nxshpDimPMww0t2uhZb5B1mRRG8utX3I0OYKH1WI6gueyPePkW1a+he6oBMll03NF9xlPHOHW6tQcHPWa6fWL2gA5v24zIK6724v2ly5cuXL8DVU2AdmK210zs2RzgHAH1E9pSSFapQcW4V6POGvIta5xr108p4Aazf7JlOlp0y5kOWQxzz8F8VfWZGTFi53Xyo5TVt1dfG5dOW83qfcU84wwMUcnM35tTRmjUWhXQGwuvB6/WDQ/udV7HeXLhF3Lgy4PguCcjZPN+XaP2UNmj5u71iYfHwHo19nOXAeFLHN2OA9WcD4LT6dnPXpETpWrZXiviBZVKjRq2wcyN7G8w0PAyvNB6JCEtW6nQeybPhbh9jsm4m4mEhSzrRVp1B0ejMkS78TbzFP1bTMCppr1GXvLl+C/BcuZQNWoJ736RhNdrHarhCibUPdb6y/IbnF3OvlcaMdnSDy3ed9IOypypwNjnpMp4v5Zebn+FxDWmu0bCyuQ8Q5lk5w5otVho216bLvL3oCccx8hrv4qfo/hqd/ZDrC/moU9weX1Zw9vuamMG0uXLly5cucUEo/Aj7yqC8f8JRVYqMDSYelt/sJryj+enLw5S7KvYiKVhNRwnl41WXMWFwRmi8JoYqIEdNRLowFrW7BlZ5l5LVd069tvHYwC9S3owt7vyCvWvqxFt/bS5cGXL8Fy5cuXAV7WrsfjyiJ021cqvhnbQ+oZpzdDmxunL6tVpfNWZU2N4UPHcrFmid4GNHPI0+GfwMYav4JQ9UDybmocOx9WnS/G5cuXLlwZcuXLhXVhd60d+vDvFUolbV3fGx0n8AlAcXGa0mCWO8BBHG+RSYin42IGvILXTtPWr3KvTwyD5VlPhJZxBvoHQo1XYg/lVv0NbC55MrAOA7wxcj9Wee77CLmXL8F3Lly5dw41bL0u6jbpWrqr42jQrfX1Xi7ek0dXHQO5V+a82clVd+Y4Pa5VLdZY4I2nIo8ObEHVogzeoft5nhFzlI0AQLult9vCyEIHcBSHFNa5x9qtnABnu7RgFUK6mCe9B+rBULNTJCD0AnnmXLly5cuXLiiYlXV1kdJ++/mUuAuCqqubFoWMKUj0UFXqtsRgFr8wGx/DK1h5OD6DbfOms5WrpxXjx8hx8RUJSZE1GFBGgAO7MFWI+YnrUMVyj0o+rrhebfUxLlsuXLly8yuCK4OzwFs/zSa3Rjx/wQXYLKacPwV1lDBN6u4vw47tDGumCPpWllV2AyyvK4ut7GNuVMJ5eF4GHXNfYCE6sA9x6J9XyDZD54ftLly5cuXFh1nPl3eM/fM/fIq6l1nomYl4YtEG+ZmjFi5jtK3BHck9Tv018F9mu7LTXBW75BEGCSYBqDQ5D3hqo++VV9CN8kjgt0IH5V8i8r6UPKa9vRRcHkY+r2VrfIcemsvwXLl+AyHnEOMLrqdYJkrWmPX+AwFwHxQnlh4MtN041fcPUnpWMKDTjnX28UjFuaxkDHalunGnfeK8EAb5Q1B5LFmBANQq14YKJe5beRynm0ecsPT0rhyvY8/rGuJ+TEm/mZ8F8ZfguXLj+PiSzxQA1nszR7L2uVenBmNWrwXB4Dc858I38LoWB6q7R592fmlOHF3ls+3eiWi4Bseesw8JRbb3F+sTWl3y81ZgOqMsUaPPV68pZG+wNXzbfrN22P0aekvwXLly5cChB1HIz9dl9yhu8totouo+CZ3sT2eI7kw12tVsNw4PkzP95zjQeF+xVyKi+ywdxSJqlQvkVg8TLXXzjbzPpcS3oUV0F56Hn9a1BEJ1MzSGAnnmXLly5cuXLlzTGdHU/wA5REGTR2TieHoP2m4m4wPMFcjUrlZYxzhJS04nJ4k9Wg/EDHGd/KsdowWpXDdD5v8AW2DzdjLly5cuXL8GGdiKqbumVSJkZR+aw7a8IvMVPEs0+YTeCeWQftMeog78INSuEw3Fh+uHTHvCXLly5cuXF60Hd5E2EGh0/wBfAw2NMGVbAetcL38ZlYwYdZBX4AftvhOrPcD9vrlHEg+8uXLl+C5eLjWHAH38/wCWPx64iXzojrqYZv4RU9kfkn3qcgjsD+frlKcU+/2ly5cuXLmXM6um/f8An6P3osHzqF3sAR6UN2OEyYSpotnDSnzX8DGb+q9Vr645Ij/ZcuXLly5w6MOh/PWfG0d/HxBE7Zbux3rhKfk95pY2VuzpL4NnyCvdZUjSh8jPrf13iLWesuXLlziSHrt/RrvjUd/Fx4lAt2mFPusNQ9l1W/ruR7NnRly5cuVnje39Gv4fbivqe3xWBkbdLLKQwsPRy+h9es7Rw+cvj4XLl3AF93+jXeHvEs8z7eNA2vaMr39pUtip6svbB9frHcwy5cuOzkB6f0Aa2jrh+0V2Yp50n38NTHo/PlLbo4OLt3dYi9sq836/q2jhly/BWOb/AEIJqb5b+kDKUO3n1LllXHpsWG6fdlK9rnnu7af+A9f2ftBhs9v6r6c9um3pFybtsw8+IbS6ODI2fl9P/AzKfM1hoVQbf1YdVw7j1JUufc59Vqatur/7PUNUT9tP3Ur/ALf8iGqOsfvifvifvp++gNRfxvil4vafp0/Tp+nT9Kn6dDUF0/qaC2fsoU5P+wKgua2yfpU/UkDozmX7QNHqP9MdyDq16tdw4cz+5jq4vRVhweMw7ArkC4uzwuSgVU7J+ififon4n6j+J+o/iZcu+ovJw/hufiIu1cYy1kbScrxPhiHwD2nxD7T4J9ofIPaHl3fXby/1VyaTpekbb4v/ALHXwaIpdFHggpLj7P4R+3pO+0BPMDfcx06y9oZOr08Rr9x4+vpvtA6AbtHtBHTwUIIzk68q7EoXnH/KX4dxwr+HlBKJiKughY6JajqDRxoz5z9p8D+0+c/aMpEJU6tdA8BUSNAs9ieob/CFCoWjR0bMcnrGwVGeYh4XGkxj8Ki5UpXzIvdtHOUOUXafLvtF3GUqC26BvOYpKuwLEraOP+MRGHU6TqOf4d+AzuFQimOt7ERjDq4TqOYqWevEV873c5ghK/ac1ld9gf8AU6+HRPlO6Gnh1hk7dTyNPRgLdUHvT6MutZqiwug2pu8to2Ch3QB9psN8N9mNj2Jwcn5ax4dpEdkjrp/ZG1AXpWhKva7lKy2NHn/sQzB9hZ8Gzxvaavm90tMpXgvghfNu+6Kun9kPI7Qc87znxzu6c4k1j+2c2U/PLV6urLd4CW0ImRGxHWCjtufseHoX2nz7YjtPDFCq0zDwHDm/eBy/hPNd+rLCqaqlul3MPLoY52odNI51xzRv1NE4wFwCdANV5ReSwy8MTTiXg24xKNM0D8EM6C0TuHOYDlTqXFTS4mXzykWS1wM7ecVqJaeVsz5jqPMhkUA3VfHmIj0/6XXy6k+D7pt4oHt9qISkJMsjy09ajr0SnkZYvw+mHzeErqsTUls65Q4mH1IPOB0Y9qjro/ZEGuJQ6WofeXiHDY6GhCEl4HTD1uFW+ENxqmVjdiSGvhV3J8O3Qp+5ifdLRV0cISnSfME+cIZntpWzb77Aq6LznPoYnIsHWhq9eESNercx2wBy0WvLyjsnqfY8PRPtPg20V/BzAK0WuA5uk0qDF4uq83MxXLhS0unkXpx8B8T22290PRO0L1w0fMHemAjeAeJrsy+UpntV2Ate0vmSnBTizKvYgo0Q8TWKKxaAS6pbTXaa2FVqMkaFFDS6lNaGIId1igPG5pVXHgbO5Hbj0yx7v/S6+PUny/Zh4s3vp0APeBxwB5p9pQTlerEAFjhHhP1z8Th9p+IQHDYKPSK+r9sVDy+yXB8YeCpPCL5+7FSc/uTF/KmGh4UdD7QyYqk6n8JSVc22yvFZQAaERDYq7vh8txI5R2/P2PD0rN6/RFfzcxoANnbHij5qWtFeak3/AInSNLvii3Wia3vKAG5h5QgakvJQfR/hikBXMvh5+FU07+Pj+3v/AGKn+ck+DbMNIoQUGt60PzMW50Q/zIbHaZPovmMkJFhu4iYaak8L8KOl9iWF8YeGC+F+S4x11PuRIm7PSGngq6X2lxcj2fwclQNyTfUw+sfUgdLNj4WRMuI8yOUVrzPZ4elZ8y2jtfncahYnqNxdrCvJLI0suza3R1H08Lg3pAdVS9S/C9pXmevASr8pf3fkmycR8LLqYpIymXws/CgorPlY/t7/ANjr49SKByDhJ+uxKwPIyzCluKNH5iQtq2rvEBKRseZKfoT/AHOUps6aHDP0D8QtYTfN9ZzxDsyjofaXH09nh6P4PiOMU2iU9dPWWjqouW4uJ5sN5W+8ddH7QoMHUck/X5+mz9ZmQCLVTcgIlNY8H8MUSnInW55SBodplCLvq92ZRX1j2eHpWdo+0y+fnwFqoc2he52zp24Rk4a+uODzJZi+AU9RY9aznnYo7kGXmgUMXh2iLUy9Ps6MJJR6y58uDvCRY01jo6kByZwafeUwVlaXRl1dY68KZ/Ey8KI5X87H/SZZiU3ac/tz9Lim10MwBXDR2/gKhKTRIHS7t3J81/Mw4jy57txVKWuqx0lFbdo5TdUcM+GuOrPKc7tzVESvn4CgDoZ+usdj1fzAPBCoHcBW2XObP9Dn6HOZN0eQqmvgRl9ef5MG75xg23Uxv7wIAccaZxdVWOXggiO8DAYgNECuiMNFH78MQJVtKQ6bnkwpyj0PvC654NeqnpNoaW9DgcDwPaLKzAq5ra9QekaaA80+5FbBqFGs53fNm/KAWWcoB+JM0w14LW3wRvLqpcquerGUYlr6qv2/tf8A1Bv/APPh/9k="

def get_logo_bytes():
    return base64.b64decode(LOGO_B64)

def get_logo_image():
    return BytesIO(get_logo_bytes())

# ─── Page Config ───
st.set_page_config(
    page_title="DentEdTech™ Evidence Engine",
    page_icon=get_logo_image(),
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─── Custom CSS ───
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=DM+Serif+Display:ital@0;1&family=DM+Sans:ital,opsz,wght@0,9..40,300;0,9..40,400;0,9..40,500;0,9..40,600;0,9..40,700;1,9..40,400&display=swap');

:root {
    --primary: #1B4D3E;
    --primary-light: #2D7A5F;
    --accent: #E8C067;
    --accent-light: #F0D48A;
    --bg-dark: #0F1A16;
    --bg-card: #162520;
    --bg-card-hover: #1C3029;
    --text-primary: #F5F8F6;
    --text-secondary: #C8D8CE;
    --text-muted: #99B3A5;
    --border: #2D4A3E;
    --danger: #E06060;
    --warning: #E8C067;
    --success: #5FCC8F;
    --video-red: #E85050;
    --recall-purple: #B48DE0;
}

.stApp {
    background-color: var(--bg-dark) !important;
    color: var(--text-primary) !important;
    font-family: 'DM Sans', sans-serif !important;
}

#MainMenu {visibility: hidden;}
footer {visibility: hidden;}

/* Show sidebar toggle on tablets/mobile (iPad) */
header[data-testid="stHeader"] {
    background: var(--bg-dark) !important;
    border-bottom: 1px solid var(--border) !important;
}

/* iPad & tablet: ensure sidebar is accessible */
@media (max-width: 1024px) {
    section[data-testid="stSidebar"] {
        z-index: 999 !important;
    }
    /* Make the hamburger menu more visible on dark bg */
    header[data-testid="stHeader"] button {
        color: var(--text-primary) !important;
    }
}

.main-header {
    text-align: center;
    padding: 2rem 1rem 1.5rem;
    border-bottom: 1px solid var(--border);
    margin-bottom: 2rem;
}
.main-header h1 {
    font-family: 'DM Serif Display', serif !important;
    font-size: 2.4rem !important;
    color: var(--text-primary) !important;
    margin: 0 !important;
    letter-spacing: -0.02em;
}
.main-header .tagline {
    font-size: 0.95rem;
    color: var(--text-secondary);
    margin-top: 0.3rem;
    font-style: italic;
}
.brand-accent { color: var(--accent) !important; }

section[data-testid="stSidebar"] {
    background-color: var(--bg-card) !important;
    border-right: 1px solid var(--border) !important;
}
section[data-testid="stSidebar"] .stMarkdown p,
section[data-testid="stSidebar"] .stMarkdown li,
section[data-testid="stSidebar"] label {
    color: var(--text-secondary) !important;
    font-family: 'DM Sans', sans-serif !important;
}
section[data-testid="stSidebar"] h1,
section[data-testid="stSidebar"] h2,
section[data-testid="stSidebar"] h3,
section[data-testid="stSidebar"] h4,
section[data-testid="stSidebar"] h5 {
    color: var(--text-primary) !important;
    font-family: 'DM Serif Display', serif !important;
}
/* Sidebar selectbox and input text */
section[data-testid="stSidebar"] .stSelectbox label,
section[data-testid="stSidebar"] .stTextInput label,
section[data-testid="stSidebar"] .stTextArea label {
    color: var(--text-secondary) !important;
}
section[data-testid="stSidebar"] .stSelectbox div[data-baseweb="select"] span,
section[data-testid="stSidebar"] .stSelectbox svg {
    color: var(--text-primary) !important;
    fill: var(--text-primary) !important;
}
/* Sidebar button text */
section[data-testid="stSidebar"] .stButton > button {
    color: var(--text-primary) !important;
}
/* Sidebar expander text */
section[data-testid="stSidebar"] .streamlit-expanderHeader p,
section[data-testid="stSidebar"] .streamlit-expanderHeader span {
    color: var(--text-primary) !important;
}

.mode-card {
    background: var(--bg-card);
    border: 1px solid var(--border);
    border-radius: 12px;
    padding: 1.5rem;
    margin-bottom: 1rem;
    transition: all 0.3s ease;
}
.mode-card:hover {
    border-color: var(--accent);
    background: var(--bg-card-hover);
}
.mode-card h3 {
    font-family: 'DM Serif Display', serif !important;
    color: var(--text-primary) !important;
    margin-top: 0 !important;
    font-size: 1.2rem !important;
}
.mode-card p {
    color: var(--text-secondary) !important;
    font-size: 0.88rem !important;
    line-height: 1.5 !important;
}

.pillar-badge {
    display: inline-block;
    padding: 3px 10px;
    border-radius: 20px;
    font-size: 0.72rem;
    font-weight: 600;
    letter-spacing: 0.04em;
    text-transform: uppercase;
    margin-right: 6px;
    margin-bottom: 4px;
}
.pillar-r { background: rgba(95, 204, 143, 0.15); color: #5FCC8F; border: 1px solid rgba(95, 204, 143, 0.3); }
.pillar-e { background: rgba(232, 192, 103, 0.15); color: #E8C067; border: 1px solid rgba(232, 192, 103, 0.3); }
.pillar-a { background: rgba(120, 170, 255, 0.15); color: #78AAFF; border: 1px solid rgba(120, 170, 255, 0.3); }
.pillar-l { background: rgba(224, 96, 96, 0.15); color: #F08080; border: 1px solid rgba(224, 96, 96, 0.3); }

.chat-msg {
    padding: 1.2rem 1.5rem;
    border-radius: 12px;
    margin-bottom: 1rem;
    line-height: 1.7;
    font-size: 0.92rem;
}
.chat-msg-user {
    background: var(--bg-card);
    border: 1px solid var(--border);
    border-left: 3px solid var(--accent);
}
.chat-msg-assistant {
    background: rgba(27, 77, 62, 0.15);
    border: 1px solid rgba(45, 122, 95, 0.25);
    border-left: 3px solid var(--primary-light);
}
.chat-msg-system {
    background: rgba(212, 168, 83, 0.08);
    border: 1px solid rgba(212, 168, 83, 0.2);
    border-left: 3px solid var(--accent);
    font-style: italic;
}

.ebl-phase {
    background: var(--bg-card);
    border: 1px solid var(--border);
    border-radius: 12px;
    padding: 1.2rem 1.5rem;
    margin-bottom: 1.5rem;
}
.ebl-phase-title {
    font-family: 'DM Serif Display', serif !important;
    font-size: 1.1rem !important;
    color: var(--accent) !important;
    margin-bottom: 0.5rem !important;
}
.ebl-phase-desc {
    color: var(--text-secondary) !important;
    font-size: 0.85rem !important;
    line-height: 1.5 !important;
}

.phase-stepper {
    display: flex;
    justify-content: space-between;
    margin-bottom: 1.5rem;
    padding: 0.8rem 0;
}
.phase-step { flex: 1; text-align: center; position: relative; padding: 0 0.5rem; }
.phase-dot {
    width: 32px; height: 32px; border-radius: 50%;
    display: inline-flex; align-items: center; justify-content: center;
    font-size: 0.75rem; font-weight: 700; margin-bottom: 0.4rem;
    transition: all 0.3s ease;
}
.phase-dot-active { background: var(--accent); color: var(--bg-dark); }
.phase-dot-done { background: var(--success); color: var(--bg-dark); }
.phase-dot-pending { background: var(--bg-card); color: var(--text-muted); border: 1px solid var(--border); }
.phase-label { font-size: 0.7rem; color: var(--text-muted); text-transform: uppercase; letter-spacing: 0.05em; }
.phase-label-active { color: var(--accent) !important; font-weight: 600; }

.reflection-box {
    background: rgba(212, 168, 83, 0.06);
    border: 1px dashed rgba(212, 168, 83, 0.35);
    border-radius: 10px;
    padding: 1.2rem 1.5rem;
    margin: 1rem 0;
}
.reflection-box h4 {
    color: var(--accent) !important;
    font-family: 'DM Serif Display', serif !important;
    font-size: 0.95rem !important;
    margin-bottom: 0.5rem !important;
}
.reflection-box p { color: var(--text-secondary) !important; font-size: 0.85rem !important; }

.limitation-notice {
    background: rgba(196, 75, 75, 0.08);
    border: 1px solid rgba(196, 75, 75, 0.2);
    border-radius: 8px;
    padding: 0.8rem 1rem;
    margin-top: 1rem;
    font-size: 0.78rem;
    color: var(--text-muted);
}
.limitation-notice strong { color: var(--danger); }

.stTextArea textarea, .stTextInput input {
    background-color: var(--bg-card) !important;
    border: 1px solid var(--border) !important;
    color: var(--text-primary) !important;
    border-radius: 10px !important;
    font-family: 'DM Sans', sans-serif !important;
}
.stTextArea textarea:focus, .stTextInput input:focus {
    border-color: var(--accent) !important;
    box-shadow: 0 0 0 1px var(--accent) !important;
}

.stButton > button {
    background-color: var(--primary) !important;
    color: var(--text-primary) !important;
    border: 1px solid var(--primary-light) !important;
    border-radius: 8px !important;
    font-family: 'DM Sans', sans-serif !important;
    font-weight: 500 !important;
    padding: 0.5rem 1.5rem !important;
    transition: all 0.2s ease !important;
}
.stButton > button:hover {
    background-color: var(--primary-light) !important;
    border-color: var(--accent) !important;
}

.stSelectbox > div > div {
    background-color: var(--bg-card) !important;
    border-color: var(--border) !important;
    color: var(--text-primary) !important;
}

.section-divider { border: none; border-top: 1px solid var(--border); margin: 1.5rem 0; }

/* ─── Chat Message Dark Theme ─── */
.stChatMessage {
    background-color: var(--bg-card) !important;
    border: 1px solid var(--border) !important;
    border-radius: 12px !important;
    color: var(--text-primary) !important;
}
.stChatMessage p, .stChatMessage li, .stChatMessage span {
    color: var(--text-primary) !important;
}
.stChatMessage h1, .stChatMessage h2, .stChatMessage h3,
.stChatMessage h4, .stChatMessage h5 {
    color: var(--text-primary) !important;
}
.stChatMessage strong {
    color: var(--accent) !important;
}
.stChatMessage a {
    color: var(--accent) !important;
}
.stChatMessage code {
    background: rgba(255,255,255,0.08) !important;
    color: var(--text-primary) !important;
    border-radius: 4px !important;
    padding: 1px 5px !important;
}
.stChatMessage pre {
    background: rgba(0,0,0,0.3) !important;
    border: 1px solid var(--border) !important;
    border-radius: 8px !important;
}

/* ─── Mobile Responsive ─── */
@media (max-width: 768px) {
    .main-header h1 {
        font-size: 1.6rem !important;
    }
    .main-header .tagline {
        font-size: 0.8rem;
    }
    .mode-card {
        padding: 1rem;
    }
    .mode-card h3 {
        font-size: 1rem !important;
    }
    .mode-card p {
        font-size: 0.82rem !important;
    }
    .pillar-badge {
        font-size: 0.65rem;
        padding: 2px 7px;
    }
    .phase-stepper {
        gap: 0.2rem;
    }
    .phase-dot {
        width: 26px;
        height: 26px;
        font-size: 0.65rem;
    }
    .phase-label {
        font-size: 0.6rem;
    }
    .ebl-phase, .recall-phase-box, .reflection-box {
        padding: 0.9rem 1rem;
    }
    .limitation-notice {
        font-size: 0.72rem;
        padding: 0.6rem 0.8rem;
    }
    .stChatMessage {
        padding: 0.8rem !important;
    }
}

.app-footer {
    text-align: center; padding: 1.5rem;
    border-top: 1px solid var(--border); margin-top: 2rem;
    color: var(--text-muted); font-size: 0.75rem;
}

/* ─── Disclaimer ─── */
.disclaimer-box {
    background: rgba(232, 192, 103, 0.06);
    border: 1px solid rgba(232, 192, 103, 0.2);
    border-radius: 10px;
    padding: 1rem 1.2rem;
    margin: 1rem 0 1.5rem;
    font-size: 0.78rem;
    color: var(--text-muted);
    line-height: 1.6;
    text-align: center;
}
.disclaimer-box strong {
    color: var(--accent);
}

/* ─── Active Recall Styles ─── */
.recall-phase-box {
    background: var(--bg-card);
    border: 1px solid var(--border);
    border-radius: 12px;
    padding: 1.2rem 1.5rem;
    margin-bottom: 1.5rem;
}
.recall-phase-title {
    font-family: 'DM Serif Display', serif;
    font-size: 1.1rem;
    color: var(--recall-purple);
    margin-bottom: 0.5rem;
}

.knowledge-bar-container { margin: 0.4rem 0; }
.knowledge-bar-label {
    font-size: 0.78rem; color: var(--text-muted);
    margin-bottom: 0.2rem; display: flex;
    justify-content: space-between;
}
.knowledge-bar-bg {
    height: 10px; background: rgba(255,255,255,0.06);
    border-radius: 5px; overflow: hidden;
}
.knowledge-bar-fill {
    height: 100%; border-radius: 5px;
    transition: width 0.6s ease;
}

.gap-missed { background: var(--danger); }
.gap-misunderstood { background: var(--warning); }
.gap-understood { background: var(--success); }

.idk-button {
    background: rgba(196, 75, 75, 0.12) !important;
    border: 1px solid rgba(196, 75, 75, 0.3) !important;
    color: #F08080 !important;
}

.recall-round-badge {
    display: inline-block;
    padding: 4px 12px;
    border-radius: 20px;
    font-size: 0.72rem;
    font-weight: 700;
    letter-spacing: 0.05em;
    text-transform: uppercase;
    background: rgba(155, 114, 207, 0.15);
    color: var(--recall-purple);
    border: 1px solid rgba(155, 114, 207, 0.3);
}
</style>
""", unsafe_allow_html=True)


# ─── System Prompts ───

EVIDENCE_SYSTEM_PROMPT = """You are the DentEdTech™ Evidence Engine, an educational AI assistant for medicine, dentistry, and pharmacology students at the University of Manchester. You operate under the REAL-AI framework principles.

## YOUR STRICT SOURCE CONSTRAINTS
You may ONLY provide information from these source types:
1. **Scientific journals** (PubMed-indexed, peer-reviewed)
2. **University websites** (.ac.uk, .edu domains)
3. **Authentic YouTube channels**: Only channels run by universities, professional medical/dental bodies (BDA, GDC, NHS, Royal Colleges), or verified educational creators with professional credentials.

You must NEVER cite Wikipedia, blogs, commercial health sites, social media, or unverified sources.

## REAL-AI FRAMEWORK INTEGRATION
### Pillar 1 — Reflective Integration
Before providing evidence, ALWAYS ask the student what they already know first. Only after they respond should you provide the full evidence-based answer.

### Pillar 3 — Authentic Clinical Alignment
Always include a **⚠️ Limitations** section. Be transparent: "This AI response is a learning aid, not clinical advice"

### Pillar 4 — Learning-Centred Partnership
Encourage the student to discuss findings with faculty. End with a reflective question.

## RESPONSE FORMAT
📋 Pre-Reflection Prompt → 🔬 Evidence Summary → 📚 Key Sources → 🎓 University Resources → 🎥 Recommended Video → ⚠️ Limitations → 🤔 Post-Learning Reflection

## CRITICAL RULES
- Never fabricate references
- Always distinguish levels of evidence
- If you cannot find strong evidence, say so honestly"""


EBL_SYSTEM_PROMPT = """You are the DentEdTech™ Enquiry-Based Learning (EBL) Facilitator. You guide students through structured inquiry WITHOUT giving direct answers. You operate under the REAL-AI framework.

## THE HYBRID EBL MODEL — 5 Phases:
1. FORMING: Encounter the problem, activate prior knowledge
2. STORMING: Generate multiple perspectives and hypotheses
3. QUESTIONING: Transform uncertainty into structured research questions (PICO/PEO)
4. SEEKING: Learn WHERE and HOW to find evidence (without providing it)
5. SYNTHESISING: Connect findings back to the original problem

## CRITICAL RULES
- NEVER provide direct evidence, citations, or links
- NEVER answer clinical questions directly
- Always respond with guiding questions
- Use "What makes you think that?" before "Have you considered...?"
- Normalise uncertainty: "Not knowing is the starting point of inquiry, not a failure"
- Always indicate current phase: 📍 **Phase [N]: [PHASE NAME]**"""


VIDEO_SEARCH_SYSTEM_PROMPT = """You are the DentEdTech™ Clinical Video Trust Engine. You find and evaluate clinical skills videos against the Video Trust Authentication Framework (VTAF) — 7 dimensions:

1. Author Credentials (25%): Degrees, postgrad, academic appointment, GDC/GMC, publications
2. Institutional Backing (20%): University/Royal College vs personal channel
3. Production Quality (10%): Multi-angle, audio, HD
4. Educational Structure (15%): Learning objectives, narration, terminology, error discussion
5. Professional Engagement (5%): Weighted LOW — niche content gets few views
6. Skill Transfer Potential (15%): Can you practise after watching?
7. Currency & Evidence (10%): Guideline alignment, publication date

Trust Levels: ≥80% ✅ TRUSTED | 60-79% ⚠️ CAUTION | <60% ❌ NOT RECOMMENDED

For each video provide: Title, URL, Channel, Author Profile (qualifications, position, registration), full 7-dimension breakdown, skill transfer assessment, and limitations.
Maximum 3 videos per query. Include direct YouTube URLs."""


ACTIVE_RECALL_ANALYSIS_PROMPT = """You are the DentEdTech™ Active Recall Analyser. You are an expert at comparing a student's recalled knowledge against their original study material to identify precise knowledge gaps.

## YOUR TASK
You will receive:
1. The ORIGINAL STUDY MATERIAL the student uploaded
2. The student's FREE RECALL attempt (what they wrote from memory) OR their ANSWERS to questions

## YOUR ANALYSIS
Compare the recall/answers against the original material and categorise EVERY key concept/fact into exactly one of three categories:

### ✅ UNDERSTOOD — Student got this right
Concepts the student recalled correctly, with accurate details and relationships.

### ⚠️ MISUNDERSTOOD — Student got this partially right or wrong
Concepts the student attempted but got details wrong, confused relationships, mixed up terminology, or had incomplete understanding. Explain exactly what they got wrong and what the correct information is.

### ❌ MISSED — Student forgot or didn't mention this
Important concepts from the study material that the student did not mention at all, or said "I don't know" to. These are complete gaps.

## RESPONSE FORMAT
You MUST respond in valid JSON with this exact structure:
```json
{
    "summary": "Brief overall assessment of the student's recall",
    "round_score": 65,
    "understood": [
        {"concept": "Name of concept", "detail": "What they got right"},
        ...
    ],
    "misunderstood": [
        {"concept": "Name of concept", "student_said": "What the student said", "correct": "What the correct information is"},
        ...
    ],
    "missed": [
        {"concept": "Name of concept", "correct": "What the student needs to learn"},
        ...
    ],
    "total_concepts": 20,
    "understood_count": 8,
    "misunderstood_count": 5,
    "missed_count": 7
}
```

## CRITICAL RULES
- Be thorough: extract EVERY important concept from the study material, not just main headings
- Include specific facts, relationships, mechanisms, definitions, clinical significance
- Be fair: if the student conveyed the right idea in different words, count it as understood
- Be precise about misunderstandings: quote what they said vs what's correct
- Do NOT include any text outside the JSON block
- round_score should be a percentage: (understood / total_concepts) × 100"""


ACTIVE_RECALL_QUESTIONS_PROMPT = """You are the DentEdTech™ Active Recall Question Generator. You generate targeted questions to test a student's knowledge, prioritising their weakest areas.

## YOUR TASK
You will receive:
1. The ORIGINAL STUDY MATERIAL
2. The student's KNOWLEDGE GAP ANALYSIS (what they understood, misunderstood, and missed)
3. The current ROUND number

## QUESTION GENERATION RULES

### Priority Order (MANDATORY):
1. FIRST: Ask about concepts the student MISSED completely (❌) — these are the biggest gaps
2. SECOND: Ask about concepts the student MISUNDERSTOOD (⚠️) — test if they now understand correctly
3. LAST: Ask about concepts they UNDERSTOOD (✅) — brief verification only

### Question Design:
- Generate 5-8 questions per round
- Questions should be specific, not vague
- Mix question types: definition, mechanism, clinical application, comparison, case-based
- For misunderstood concepts: frame questions that specifically target the misconception
- For missed concepts: start with foundational questions before complex ones

## RESPONSE FORMAT
You MUST respond in valid JSON:
```json
{
    "questions": [
        {
            "id": 1,
            "question": "The question text",
            "concept": "Which concept this tests",
            "gap_type": "missed",
            "difficulty": "foundation"
        },
        ...
    ],
    "focus_message": "Brief message to the student about what this round focuses on"
}
```

gap_type must be one of: "missed", "misunderstood", "understood"
difficulty must be one of: "foundation", "application", "integration"

## CRITICAL RULES
- At least 60% of questions should target missed or misunderstood concepts
- Never provide answers in the questions
- Each question should test ONE concept clearly
- Include "I don't know" as a valid response option — tell the student this in focus_message
- Do NOT include any text outside the JSON block"""


ACTIVE_RECALL_RELEARN_PROMPT = """You are the DentEdTech™ Active Recall Re-Learning Presenter. You re-present study material in a prioritised order based on the student's knowledge gaps.

## YOUR TASK
You will receive:
1. The ORIGINAL STUDY MATERIAL
2. The student's KNOWLEDGE GAP ANALYSIS

## RE-PRESENTATION ORDER (MANDATORY):
Present the material in this EXACT priority order:

### 1. ❌ MISSED CONCEPTS — Present FIRST
These are concepts the student completely forgot or never learned. Present each one with:
- Clear definition/explanation
- Why it matters clinically
- Memory hook or mnemonic if helpful
- Connection to concepts they already understand

### 2. ⚠️ MISUNDERSTOOD CONCEPTS — Present SECOND
These are concepts the student got wrong. For each:
- State what they incorrectly believed
- Explain why that's wrong
- Present the correct information clearly
- Highlight the specific distinction they missed

### 3. ✅ UNDERSTOOD CONCEPTS — Present LAST (brief)
These concepts the student already knows. Present briefly as confirmation/reinforcement.
- Brief summary only
- Any nuances they might deepen

## FORMAT
Use clear headings and structure. Make it scannable. Use clinical examples from the original material.
Start with: "Here's your study material, reorganised based on your recall performance. We're starting with what needs the most attention."

End with: "When you're ready, we'll test you again on the areas you struggled with. Take your time reading through — especially the ❌ sections."

## CRITICAL RULES
- Use the EXACT content from the original study material — don't invent new information
- Be encouraging, not punitive
- Emphasise that forgetting is normal and part of the learning process
- Keep the re-presentation focused and scannable"""


ACTIVE_RECALL_GRADE_ANSWERS_PROMPT = """You are the DentEdTech™ Active Recall Answer Grader. You grade a student's answers to targeted questions by comparing them against the original study material.

## YOUR TASK
You will receive:
1. The ORIGINAL STUDY MATERIAL
2. The QUESTIONS that were asked
3. The STUDENT'S ANSWERS to each question
4. The INITIAL GAP ANALYSIS from their free recall (so you know their starting point)

## GRADING RULES
For each question-answer pair, provide:
- Whether the answer is CORRECT, PARTIALLY CORRECT, INCORRECT, or the student said "I DON'T KNOW"
- The correct answer from the study material
- A brief explanation of what was right or wrong
- Whether this represents an IMPROVEMENT from their initial recall (they got it wrong before but right now), CONSISTENT (same as before), or REGRESSION (they knew it before but got it wrong now)

## RESPONSE FORMAT
You MUST respond in valid JSON:
```json
{
    "graded_answers": [
        {
            "question_id": 1,
            "question": "The question text",
            "student_answer": "What the student wrote",
            "verdict": "correct",
            "correct_answer": "The correct answer from the study material",
            "explanation": "Brief explanation of why this is correct/incorrect",
            "improvement_status": "improved",
            "concept": "Name of concept tested"
        }
    ],
    "questions_correct": 5,
    "questions_partial": 1,
    "questions_incorrect": 1,
    "questions_idk": 1,
    "total_questions": 8,
    "overall_feedback": "Brief encouraging summary of performance on these questions"
}
```

verdict must be one of: "correct", "partially_correct", "incorrect", "i_dont_know"
improvement_status must be one of: "improved", "consistent", "regression", "new_knowledge"

## CRITICAL RULES
- Grade fairly: if the student conveys the right idea in different words, mark as correct
- Be specific about what's wrong in partially correct answers
- For "I DON'T KNOW" answers, don't penalise — it's honest and shows the student knows their limits
- Always provide the correct answer so the student can learn from each question
- Do NOT include any text outside the JSON block"""


# ─── Trusted Channel Registry ───
TRUSTED_CHANNELS = {
    "university": [
        {"channel": "University of Manchester", "url": "https://www.youtube.com/@OfficialUoM", "category": "University", "trust_floor": 85, "notes": "Home institution."},
        {"channel": "King's College London Dentistry", "url": "https://www.youtube.com/@KCLDentistry", "category": "University", "trust_floor": 90, "notes": "Leading UK dental school."},
        {"channel": "Harvard School of Dental Medicine", "url": "https://www.youtube.com/@HarvardDentalMedicine", "category": "University", "trust_floor": 90, "notes": "International leader."},
        {"channel": "University of Michigan School of Dentistry", "url": "https://www.youtube.com/@umichdent", "category": "University", "trust_floor": 88, "notes": "Extensive clinical skills library."},
    ],
    "professional_bodies": [
        {"channel": "British Dental Association (BDA)", "url": "https://www.youtube.com/@TheBDA", "category": "Professional Body", "trust_floor": 90, "notes": "UK professional body."},
        {"channel": "General Dental Council (GDC)", "url": "https://www.youtube.com/@TheGDCUK", "category": "Regulator", "trust_floor": 85, "notes": "UK dental regulator."},
        {"channel": "Royal College of Surgeons of England", "url": "https://www.youtube.com/@RCSEngland", "category": "Royal College", "trust_floor": 92, "notes": "FDSRCS, surgical technique."},
        {"channel": "Royal Pharmaceutical Society", "url": "https://www.youtube.com/@royalpharmaceuticalsociety", "category": "Professional Body", "trust_floor": 88, "notes": "Pharmacology content."},
    ],
    "nhs": [
        {"channel": "NHS England", "url": "https://www.youtube.com/@NHSEngland", "category": "NHS", "trust_floor": 85, "notes": "Official NHS channel."},
    ],
}


# ─── Session State ───
def init_session():
    defaults = {
        "mode": None,
        "evidence_messages": [],
        "ebl_messages": [],
        "ebl_phase": 1,
        "ebl_case": None,
        "reflection_given": False,
        "discipline": "Dentistry",
        "year_of_study": "Year 3",
        "video_messages": [],
        # Active Recall state
        "ar_phase": "upload",  # upload → free_recall → analysis → questions → relearn → repeat
        "ar_study_material": None,
        "ar_file_name": None,
        "ar_free_recall": None,
        "ar_analysis": None,
        "ar_questions": None,
        "ar_answers": {},
        "ar_round": 1,
        "ar_history": [],
        "ar_messages": [],
        # Feedback
        "feedback_submitted": False,
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v

init_session()


# ─── Helper Functions ───

def get_api_key():
    try:
        return st.secrets["ANTHROPIC_API_KEY"]
    except (KeyError, FileNotFoundError):
        return None


def log_user_session():
    """Log each unique session to Google Sheets for user counting. Runs once per session."""
    if st.session_state.get("session_logged"):
        return  # Already logged this session

    try:
        import gspread
        from google.oauth2.service_account import Credentials

        scopes = [
            "https://www.googleapis.com/auth/spreadsheets",
            "https://www.googleapis.com/auth/drive",
        ]
        creds_dict = dict(st.secrets["gcp_service_account"])
        creds = Credentials.from_service_account_info(creds_dict, scopes=scopes)
        client = gspread.authorize(creds)

        sheet = client.open("DentEdTech_UserLog").sheet1

        # Generate a unique session ID
        session_id = str(uuid.uuid4())[:8]
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        # Append row: timestamp, session_id, discipline, year
        sheet.append_row([
            timestamp,
            session_id,
            st.session_state.get("discipline", "Unknown"),
            st.session_state.get("year_of_study", "Unknown"),
        ])

        st.session_state.session_logged = True
    except Exception:
        # Silently fail — don't break the app if tracking fails
        st.session_state.session_logged = True


def get_user_count():
    """Get total session count from Google Sheets."""
    try:
        import gspread
        from google.oauth2.service_account import Credentials

        scopes = [
            "https://www.googleapis.com/auth/spreadsheets",
            "https://www.googleapis.com/auth/drive",
        ]
        creds_dict = dict(st.secrets["gcp_service_account"])
        creds = Credentials.from_service_account_info(creds_dict, scopes=scopes)
        client = gspread.authorize(creds)

        sheet = client.open("DentEdTech_UserLog").sheet1
        # Row count minus header row
        return max(0, len(sheet.get_all_values()) - 1)
    except Exception:
        return None


def call_claude(messages, system_prompt, use_search=False):
    api_key = get_api_key()
    if not api_key:
        return "⚠️ API key not found. Please add ANTHROPIC_API_KEY to your Streamlit secrets."

    client = anthropic.Anthropic(api_key=api_key)
    kwargs = {
        "model": "claude-sonnet-4-20250514",
        "max_tokens": 4096,
        "system": system_prompt,
        "messages": messages,
    }
    if use_search:
        kwargs["tools"] = [{"type": "web_search_20250305", "name": "web_search"}]

    try:
        response = client.messages.create(**kwargs)
        text_parts = []
        for block in response.content:
            if hasattr(block, "text"):
                text_parts.append(block.text)
        return "\n".join(text_parts) if text_parts else "I wasn't able to generate a response. Please try again."
    except anthropic.AuthenticationError:
        return "⚠️ Invalid API key."
    except anthropic.RateLimitError:
        return "⚠️ Rate limit reached. Please wait."
    except Exception as e:
        return f"⚠️ An error occurred: {str(e)}"


def call_claude_json(messages, system_prompt):
    """Call Claude and parse JSON from response."""
    raw = call_claude(messages, system_prompt, use_search=False)
    # Extract JSON from response (handle markdown code blocks)
    json_match = re.search(r'```(?:json)?\s*([\s\S]*?)```', raw)
    if json_match:
        raw = json_match.group(1)
    # Try to find JSON object or array
    json_match = re.search(r'(\{[\s\S]*\})', raw)
    if json_match:
        try:
            return json.loads(json_match.group(1))
        except json.JSONDecodeError:
            return None
    return None


def render_phase_stepper(current_phase):
    phases = [("1", "Forming"), ("2", "Storming"), ("3", "Questioning"), ("4", "Seeking"), ("5", "Synthesising")]
    html = '<div class="phase-stepper">'
    for num, label in phases:
        phase_num = int(num)
        if phase_num < current_phase:
            dot_class, label_class, dot_content = "phase-dot phase-dot-done", "phase-label", "✓"
        elif phase_num == current_phase:
            dot_class, label_class, dot_content = "phase-dot phase-dot-active", "phase-label phase-label-active", num
        else:
            dot_class, label_class, dot_content = "phase-dot phase-dot-pending", "phase-label", num
        html += f'<div class="phase-step"><div class="{dot_class}">{dot_content}</div><div class="{label_class}">{label}</div></div>'
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)


def render_message(role, content):
    """Render chat message using Streamlit's native chat_message for proper markdown rendering."""
    if role == "user":
        with st.chat_message("user", avatar="🧑‍🎓"):
            st.markdown(content)
    elif role == "system":
        with st.chat_message("assistant", avatar="🔔"):
            st.markdown(f"*{content}*")
    else:
        with st.chat_message("assistant", avatar="🔬"):
            st.markdown(content)


def render_real_ai_badges(pillars):
    badge_map = {
        "R": ("pillar-r", "Reflective Integration"), "E": ("pillar-e", "Equity by Design"),
        "A": ("pillar-a", "Authentic Alignment"), "L": ("pillar-l", "Learning Partnership"),
    }
    return "".join(f'<span class="pillar-badge {badge_map[p][0]}">{badge_map[p][1]}</span>' for p in pillars)


def render_youtube_embed(video_id):
    return f"""<div style="position:relative;padding-bottom:56.25%;height:0;overflow:hidden;border-radius:10px;margin:1rem 0;">
        <iframe src="https://www.youtube.com/embed/{video_id}" style="position:absolute;top:0;left:0;width:100%;height:100%;border:none;border-radius:10px;"
        allow="accelerometer;autoplay;clipboard-write;encrypted-media;gyroscope;picture-in-picture" allowfullscreen></iframe></div>"""


def render_knowledge_bar(label, count, total, color_class):
    pct = (count / total * 100) if total > 0 else 0
    return f"""<div class="knowledge-bar-container">
        <div class="knowledge-bar-label"><span>{label}</span><span>{count}/{total} ({pct:.0f}%)</span></div>
        <div class="knowledge-bar-bg"><div class="knowledge-bar-fill {color_class}" style="width:{pct}%"></div></div></div>"""


def render_recall_phase_stepper(current_phase):
    phases = [("upload", "📄 Upload"), ("free_recall", "✍️ Recall"), ("gap_report_1", "📊 Gaps"),
              ("questions", "❓ Questions"), ("q_feedback", "📝 Feedback"), ("gap_report_2", "📊 Final")]
    html = '<div class="phase-stepper">'
    phase_order = ["upload", "free_recall", "gap_report_1", "questions", "q_feedback", "gap_report_2"]
    current_idx = phase_order.index(current_phase) if current_phase in phase_order else 0
    for i, (key, label) in enumerate(phases):
        if i < current_idx:
            dot_class, label_class = "phase-dot phase-dot-done", "phase-label"
            dot_content = "✓"
        elif i == current_idx:
            dot_class, label_class = "phase-dot phase-dot-active", "phase-label phase-label-active"
            dot_content = str(i + 1)
        else:
            dot_class, label_class = "phase-dot phase-dot-pending", "phase-label"
            dot_content = str(i + 1)
        html += f'<div class="phase-step"><div class="{dot_class}">{dot_content}</div><div class="{label_class}">{label}</div></div>'
    html += "</div>"
    st.markdown(html, unsafe_allow_html=True)


def read_uploaded_file(uploaded_file):
    """Extract text from uploaded file."""
    if uploaded_file is None:
        return None
    name = uploaded_file.name.lower()
    try:
        if name.endswith('.txt') or name.endswith('.md'):
            return uploaded_file.read().decode('utf-8')
        elif name.endswith('.pdf'):
            import pypdf
            reader = pypdf.PdfReader(uploaded_file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text
        elif name.endswith('.docx'):
            import docx
            doc = docx.Document(uploaded_file)
            return "\n".join([para.text for para in doc.paragraphs])
        elif name.endswith('.pptx'):
            from pptx import Presentation
            prs = Presentation(uploaded_file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text + "\n"
            return text
        else:
            return uploaded_file.read().decode('utf-8', errors='ignore')
    except Exception as e:
        return f"Error reading file: {str(e)}"


def _render_gap_report(analysis, context_label):
    """Render the gap analysis report — used by both gap_report_1 and gap_report_2."""
    a = analysis
    total = a.get("total_concepts", 0)
    understood = a.get("understood_count", 0)
    misunderstood = a.get("misunderstood_count", 0)
    missed = a.get("missed_count", 0)
    score = a.get("round_score", 0)

    if score >= 80:
        score_color = "#5FCC8F"
    elif score >= 60:
        score_color = "#E8C067"
    else:
        score_color = "#E06060"

    r, g, b = int(score_color[1:3], 16), int(score_color[3:5], 16), int(score_color[5:7], 16)
    st.markdown(f"""<div style="text-align:center;margin:1.5rem 0;">
        <div style="display:inline-flex;width:90px;height:90px;border-radius:50%;border:4px solid {score_color};
        background:rgba({r},{g},{b},0.12);align-items:center;justify-content:center;
        font-family:'DM Serif Display',serif;font-size:1.6rem;color:{score_color};">{score}%</div>
        <div style="color:var(--text-secondary);margin-top:0.5rem;font-size:0.85rem;"><em>{context_label}</em></div>
        <div style="color:var(--text-muted);margin-top:0.2rem;font-size:0.82rem;">{a.get('summary', '')}</div></div>""", unsafe_allow_html=True)

    st.markdown(render_knowledge_bar("✅ Understood", understood, total, "gap-understood"), unsafe_allow_html=True)
    st.markdown(render_knowledge_bar("⚠️ Misunderstood", misunderstood, total, "gap-misunderstood"), unsafe_allow_html=True)
    st.markdown(render_knowledge_bar("❌ Missed", missed, total, "gap-missed"), unsafe_allow_html=True)

    with st.expander("❌ Missed Concepts — What you need to learn", expanded=True):
        items = a.get("missed", [])
        if items:
            for item in items:
                st.markdown(f"**{item['concept']}**")
                st.markdown(f"{item['correct']}")
                st.markdown("---")
        else:
            st.markdown("*Nothing missed — well done!*")

    with st.expander("⚠️ Misunderstood Concepts — What you got wrong", expanded=True):
        items = a.get("misunderstood", [])
        if items:
            for item in items:
                st.markdown(f"**{item['concept']}**")
                st.markdown(f"You said: *\"{item['student_said']}\"*")
                st.markdown(f"Correct: **{item['correct']}**")
                st.markdown("---")
        else:
            st.markdown("*No misunderstandings — well done!*")

    with st.expander("✅ Understood Concepts — What you got right", expanded=False):
        items = a.get("understood", [])
        if items:
            for item in items:
                st.markdown(f"**{item['concept']}** — {item['detail']}")
        else:
            st.markdown("*Keep working — you'll get there!*")


# ─── Log User Session ───
log_user_session()

# ─── Sidebar ───
with st.sidebar:
    st.markdown('<div style="padding: 0.5rem 3rem 0;">', unsafe_allow_html=True)
    st.image(get_logo_bytes(), use_container_width=True)
    st.markdown('</div>', unsafe_allow_html=True)
    st.markdown("""
    <div style="text-align:center; padding: 0 0 0.5rem;">
        <div style="font-size: 0.72rem; color: #99B3A5; letter-spacing: 0.08em; text-transform: uppercase;">
            Evidence Engine
        </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("<hr class='section-divider'>", unsafe_allow_html=True)

    st.markdown("##### 🎓 Your Profile")
    st.session_state.discipline = st.selectbox("Discipline", ["Dentistry", "Medicine", "Pharmacology"], index=0)
    st.session_state.year_of_study = st.selectbox("Year of Study", ["Year 1", "Year 2", "Year 3", "Year 4", "Year 5", "Postgraduate"], index=2)

    st.markdown("<hr class='section-divider'>", unsafe_allow_html=True)

    with st.expander("📐 About the REAL-AI Framework"):
        st.markdown("""
        **R** — Reflective Integration · **E** — Equity by Design
        **A** — Authentic Clinical Alignment · **L** — Learning-Centred Partnership
        *Framework: Beyond the Algorithm (2026)*
        """)

    st.markdown("<hr class='section-divider'>", unsafe_allow_html=True)

    if st.session_state.mode is not None:
        if st.button("← Back to Mode Selection", use_container_width=True, key="btn_back"):
            # Clear state for current mode before going back
            if st.session_state.mode == "evidence":
                st.session_state.evidence_messages = []
                st.session_state.reflection_given = False
            elif st.session_state.mode == "ebl":
                st.session_state.ebl_messages = []
                st.session_state.ebl_phase = 1
                st.session_state.ebl_case = None
            elif st.session_state.mode == "video":
                st.session_state.video_messages = []
            elif st.session_state.mode == "recall":
                for k in ["ar_phase", "ar_study_material", "ar_file_name", "ar_free_recall",
                           "ar_analysis", "ar_analysis_post_questions", "ar_question_answers",
                           "ar_graded_answers", "ar_questions", "ar_answers", "ar_round",
                           "ar_history", "ar_messages", "ar_relearn_content"]:
                    if k in st.session_state:
                        del st.session_state[k]
            st.session_state.mode = None
            st.rerun()
        if st.button("🔄 Reset Conversation", use_container_width=True, key="btn_reset"):
            if st.session_state.mode == "evidence":
                st.session_state.evidence_messages = []
                st.session_state.reflection_given = False
            elif st.session_state.mode == "ebl":
                st.session_state.ebl_messages = []
                st.session_state.ebl_phase = 1
                st.session_state.ebl_case = None
            elif st.session_state.mode == "video":
                st.session_state.video_messages = []
            elif st.session_state.mode == "recall":
                for k in ["ar_phase", "ar_study_material", "ar_file_name", "ar_free_recall",
                           "ar_analysis", "ar_analysis_post_questions", "ar_question_answers",
                           "ar_graded_answers", "ar_questions", "ar_answers", "ar_round",
                           "ar_history", "ar_messages", "ar_relearn_content"]:
                    if k in st.session_state:
                        del st.session_state[k]
                init_session()
            st.rerun()

    st.markdown("<hr class='section-divider'>", unsafe_allow_html=True)

    # Feedback button — always visible
    if st.button("💬 Give Feedback", use_container_width=True, key="btn_feedback_sidebar"):
        st.session_state.mode = "feedback"
        st.session_state.feedback_submitted = False
        st.rerun()

    st.markdown("<hr class='section-divider'>", unsafe_allow_html=True)

    # User count
    user_count = get_user_count()
    if user_count is not None:
        st.markdown(f"""<div style="text-align:center;padding:0.5rem 0;">
            <span style="font-size:1.3rem;font-weight:700;color:var(--accent);">{user_count}</span>
            <span style="font-size:0.75rem;color:var(--text-muted);display:block;">sessions to date</span>
        </div>""", unsafe_allow_html=True)

    st.markdown("""<div class="app-footer">
        <strong style="color:var(--accent);font-size:0.7rem;">⚠️ PROTOTYPE — NOT AN OFFICIAL UNIVERSITY APPLICATION</strong><br>
        <span style="font-size:0.68rem;">Verify all outputs against your official curriculum.</span><br><br>
        © 2026 DentEdTech™<br>University of Manchester<br><em>Not a substitute for clinical judgement</em></div>""", unsafe_allow_html=True)


# ─── Main Content ───
st.markdown("""
<div class="main-header">
    <h1>Dent<span class="brand-accent">Ed</span>Tech™ Evidence Engine</h1>
    <div class="tagline">Theory-informed AI for health professions learning — built on the REAL-AI framework</div>
</div>
""", unsafe_allow_html=True)


# ─── Mode Selection ───
if st.session_state.mode is None:
    st.markdown(f"""<div style="text-align:center;margin-bottom:2rem;">
        <span style="color:var(--text-secondary);font-size:0.9rem;">Welcome, {st.session_state.discipline} student · {st.session_state.year_of_study} · University of Manchester</span></div>""", unsafe_allow_html=True)

    st.markdown("""<div class="disclaimer-box">
        <strong>⚠️ Important Notice</strong><br>
        This platform is <strong>not an official application</strong> of the University of Manchester or any other university.
        It is a <strong>research prototype</strong> developed under the DentEdTech™ project.
        All outputs must be verified against your official curriculum, lecture materials, and clinical guidelines.
        Always consult your tutors and supervisors before applying any information in clinical practice.
    </div>""", unsafe_allow_html=True)

    col1, col2 = st.columns(2, gap="medium")
    col3, col4 = st.columns(2, gap="medium")

    with col1:
        st.markdown(f"""<div class="mode-card"><h3>🔬 Evidence-Based Knowledge</h3>
            <p>Ask clinical or scientific questions and receive evidence-based answers sourced exclusively from peer-reviewed journals, university resources, and verified educational videos.</p>
            <div style="margin-top:0.8rem;">{render_real_ai_badges(["R", "A"])}</div></div>""", unsafe_allow_html=True)
        if st.button("Enter Evidence Mode →", key="btn_evidence", use_container_width=True):
            st.session_state.mode = "evidence"
            st.rerun()

    with col2:
        st.markdown(f"""<div class="mode-card"><h3>🧭 Enquiry-Based Learning</h3>
            <p>Develop your inquiry skills through a guided 5-phase cycle: Forming, Storming, Questioning, Seeking, and Synthesising. The AI guides you to discover answers yourself.</p>
            <div style="margin-top:0.8rem;">{render_real_ai_badges(["R", "E", "A", "L"])}</div></div>""", unsafe_allow_html=True)
        if st.button("Enter EBL Mode →", key="btn_ebl", use_container_width=True):
            st.session_state.mode = "ebl"
            st.session_state.ebl_messages = [{"role": "assistant", "content": "📍 **Phase 1: FORMING**\n\nWelcome to Enquiry-Based Learning.\n\nYou can either:\n- **Bring your own case** — describe a clinical scenario or problem\n- **Ask me for a case** — tell me the subject area\n\n*What topic are you most curious about right now?*"}]
            st.rerun()

    with col3:
        st.markdown(f"""<div class="mode-card"><h3>🎥 Clinical Video Trust Engine</h3>
            <p>Find the most trustworthy clinical skills videos on YouTube, scored against a 7-dimension trust framework. Videos play directly in the platform with full trust breakdowns.</p>
            <div style="margin-top:0.8rem;">{render_real_ai_badges(["R", "A", "L"])}</div></div>""", unsafe_allow_html=True)
        if st.button("Enter Video Mode →", key="btn_video", use_container_width=True):
            st.session_state.mode = "video"
            st.rerun()

    with col4:
        st.markdown(f"""<div class="mode-card"><h3>🧠 Active Recall</h3>
            <p>Upload your study material, write everything you remember, then let AI identify what you understood, misunderstood, and missed completely. Material is re-presented starting with your biggest gaps. Repeat until mastery.</p>
            <div style="margin-top:0.8rem;">{render_real_ai_badges(["R", "L"])}</div></div>""", unsafe_allow_html=True)
        if st.button("Enter Active Recall →", key="btn_recall", use_container_width=True):
            st.session_state.mode = "recall"
            st.rerun()


# ─── Evidence-Based Mode ───
elif st.session_state.mode == "evidence":
    st.markdown(f"""<div style="margin-bottom:1.5rem;"><span style="font-family:'DM Serif Display',serif;font-size:1.4rem;color:var(--text-primary);">🔬 Evidence-Based Knowledge</span>
        <span style="margin-left:1rem;">{render_real_ai_badges(["R", "A"])}</span></div>""", unsafe_allow_html=True)
    st.markdown("""<div class="limitation-notice"><strong>⚠️ Pillar 3 — Transparency:</strong> This AI searches peer-reviewed sources. Always verify against primary sources and discuss with supervisors.</div>""", unsafe_allow_html=True)
    st.markdown("<br>", unsafe_allow_html=True)
    for msg in st.session_state.evidence_messages:
        render_message(msg["role"], msg["content"])
    user_input = st.chat_input("Ask a clinical or scientific question...", key="evidence_input")
    if user_input:
        st.session_state.evidence_messages.append({"role": "user", "content": user_input})
        context = f"[Student: {st.session_state.discipline}, {st.session_state.year_of_study}, University of Manchester]"
        api_msgs = [{"role": m["role"], "content": (f"{context}\n\n{m['content']}" if i == 0 and m["role"] == "user" else m["content"])} for i, m in enumerate(st.session_state.evidence_messages)]
        with st.spinner("Searching evidence-based sources..."):
            response = call_claude(api_msgs, EVIDENCE_SYSTEM_PROMPT, use_search=True)
        st.session_state.evidence_messages.append({"role": "assistant", "content": response})
        st.rerun()


# ─── EBL Mode ───
elif st.session_state.mode == "ebl":
    st.markdown(f"""<div style="margin-bottom:1rem;"><span style="font-family:'DM Serif Display',serif;font-size:1.4rem;color:var(--text-primary);">🧭 Enquiry-Based Learning</span>
        <span style="margin-left:1rem;">{render_real_ai_badges(["R", "E", "A", "L"])}</span></div>""", unsafe_allow_html=True)
    render_phase_stepper(st.session_state.ebl_phase)
    phase_info = {1: ("Forming", "Encounter the problem. Activate what you already know."), 2: ("Storming", "Generate hypotheses freely. Challenge assumptions."),
                  3: ("Questioning", "Transform uncertainty into structured research questions."), 4: ("Seeking", "Plan your evidence search strategy."), 5: ("Synthesising", "Connect findings back to the case.")}
    pn, pd = phase_info[st.session_state.ebl_phase]
    st.markdown(f'<div class="ebl-phase"><div class="ebl-phase-title">📍 Phase {st.session_state.ebl_phase}: {pn}</div><div class="ebl-phase-desc">{pd}</div></div>', unsafe_allow_html=True)
    st.markdown("""<div class="limitation-notice"><strong>⚠️ EBL Commitment:</strong> This mode will NOT give you direct answers. The struggle of finding answers yourself is where deep learning happens.</div>""", unsafe_allow_html=True)
    st.markdown("<br>", unsafe_allow_html=True)
    for msg in st.session_state.ebl_messages:
        render_message(msg["role"], msg["content"])
    nc1, nc2, nc3 = st.columns([1, 2, 1])
    with nc1:
        if st.session_state.ebl_phase > 1 and st.button("← Previous Phase"):
            st.session_state.ebl_phase -= 1
            st.session_state.ebl_messages.append({"role": "assistant", "content": f"📍 Moving back to **Phase {st.session_state.ebl_phase}: {phase_info[st.session_state.ebl_phase][0]}**"})
            st.rerun()
    with nc3:
        if st.session_state.ebl_phase < 5 and st.button("Next Phase →"):
            st.session_state.ebl_phase += 1
            st.session_state.ebl_messages.append({"role": "assistant", "content": f"📍 Progressing to **Phase {st.session_state.ebl_phase}: {phase_info[st.session_state.ebl_phase][0]}**\n\n{phase_info[st.session_state.ebl_phase][1]}"})
            st.rerun()
    user_input = st.chat_input("Share your thinking...", key="ebl_input")
    if user_input:
        st.session_state.ebl_messages.append({"role": "user", "content": user_input})
        ctx = f"[Student: {st.session_state.discipline}, {st.session_state.year_of_study}]\n[Phase: {st.session_state.ebl_phase} — {phase_info[st.session_state.ebl_phase][0]}]"
        api_msgs = []
        for i, m in enumerate(st.session_state.ebl_messages):
            if m["role"] == "user" and i == len(st.session_state.ebl_messages) - 1:
                api_msgs.append({"role": "user", "content": f"{ctx}\n\n{m['content']}"})
            else:
                api_msgs.append(m)
        with st.spinner("Reflecting on your inquiry..."):
            response = call_claude(api_msgs, EBL_SYSTEM_PROMPT, use_search=False)
        if st.session_state.ebl_phase < 5:
            trans = {1: ["phase 2", "move to storming"], 2: ["phase 3", "move to questioning"], 3: ["phase 4", "move to seeking"], 4: ["phase 5", "move to synthesising"]}
            if any(kw in response.lower() for kw in trans.get(st.session_state.ebl_phase, [])):
                st.session_state.ebl_phase += 1
        st.session_state.ebl_messages.append({"role": "assistant", "content": response})
        st.rerun()


# ─── Video Trust Engine ───
elif st.session_state.mode == "video":
    st.markdown(f"""<div style="margin-bottom:1rem;"><span style="font-family:'DM Serif Display',serif;font-size:1.4rem;color:var(--text-primary);">🎥 Clinical Video Trust Engine</span>
        <span style="margin-left:1rem;">{render_real_ai_badges(["R", "A", "L"])}</span></div>""", unsafe_allow_html=True)
    st.markdown("""<div class="ebl-phase"><div class="ebl-phase-title">📐 Video Trust Authentication Framework (VTAF)</div>
        <div class="ebl-phase-desc">7 dimensions: Author Credentials (25%), Institutional Backing (20%), Educational Structure (15%), Skill Transfer (15%), Production Quality (10%), Currency (10%), Engagement (5%). Only ≥60% recommended.</div></div>""", unsafe_allow_html=True)
    st.markdown("""<div class="limitation-notice"><strong>⚠️ Transparency:</strong> Trust scores are AI-assessed. Verify author credentials independently. Watching does not replace supervised practice.</div>""", unsafe_allow_html=True)
    st.markdown("<br>", unsafe_allow_html=True)
    if not st.session_state.video_messages:
        st.markdown("""<div class="reflection-box"><h4>🤔 Before You Search (Pillar 1)</h4><p>What do you already know about this procedure? What specific aspect are you uncertain about?</p></div>""", unsafe_allow_html=True)
    for msg in st.session_state.video_messages:
        render_message(msg["role"], msg["content"])
        if msg["role"] == "assistant":
            for vid_id in dict.fromkeys(re.findall(r'(?:youtube\.com/watch\?v=|youtu\.be/)([\w-]+)', msg["content"])):
                st.markdown(render_youtube_embed(vid_id), unsafe_allow_html=True)
    user_input = st.chat_input("Search for a clinical skill or procedure...", key="video_input")
    if user_input:
        st.session_state.video_messages.append({"role": "user", "content": user_input})
        trusted_ctx = "\n\n## PRE-VERIFIED CHANNELS\n" + "".join(f"- {ch['channel']} ({ch['category']}) — Floor: {ch['trust_floor']}%\n" for cat in TRUSTED_CHANNELS.values() for ch in cat)
        ctx = f"[Student: {st.session_state.discipline}, {st.session_state.year_of_study}]\n[Search YouTube. Evaluate against VTAF. Provide direct URLs.]"
        api_msgs = [{"role": m["role"], "content": (f"{ctx}\n\n{m['content']}" if m["role"] == "user" and i == len(st.session_state.video_messages) - 1 else m["content"])} for i, m in enumerate(st.session_state.video_messages)]
        with st.spinner("Searching and evaluating clinical videos..."):
            response = call_claude(api_msgs, VIDEO_SEARCH_SYSTEM_PROMPT + trusted_ctx, use_search=True)
        st.session_state.video_messages.append({"role": "assistant", "content": response})
        st.rerun()


# ─── Active Recall Mode ───
elif st.session_state.mode == "recall":

    st.markdown(f"""<div style="margin-bottom:1rem;"><span style="font-family:'DM Serif Display',serif;font-size:1.4rem;color:var(--text-primary);">
        🧠 Active Recall</span><span style="margin-left:1rem;">{render_real_ai_badges(["R", "L"])}</span></div>""", unsafe_allow_html=True)

    # Round indicator
    st.markdown(f'<span class="recall-round-badge">Round {st.session_state.ar_round}</span>', unsafe_allow_html=True)
    st.markdown("<br>", unsafe_allow_html=True)

    # Phase stepper
    render_recall_phase_stepper(st.session_state.ar_phase)

    # ── PHASE 1: Upload ──
    if st.session_state.ar_phase == "upload":
        st.markdown("""<div class="recall-phase-box"><div class="recall-phase-title">📄 Step 1: Upload Your Study Material</div>
            <div class="ebl-phase-desc">Upload the material you've been studying. This can be lecture notes, a textbook chapter, a PDF, Word document, PowerPoint, or plain text file. The AI will use this as the source of truth to assess your recall.</div></div>""", unsafe_allow_html=True)

        st.markdown("""<div class="reflection-box"><h4>🤔 Pillar 1: Reflective Integration</h4>
            <p>Before you begin, close your notes. The power of active recall comes from retrieving information without looking. After uploading, you'll write everything you remember — no peeking.</p></div>""", unsafe_allow_html=True)

        uploaded_file = st.file_uploader(
            "Upload your study material",
            type=["txt", "md", "pdf", "docx", "pptx"],
            help="Supported: PDF, Word, PowerPoint, Text files",
        )

        if uploaded_file is not None:
            with st.spinner("Reading your study material..."):
                content = read_uploaded_file(uploaded_file)
            if content and not content.startswith("Error"):
                st.session_state.ar_study_material = content
                st.session_state.ar_file_name = uploaded_file.name
                st.success(f"✅ Loaded: **{uploaded_file.name}** ({len(content.split())} words)")
                if st.button("I've read this material — let's test my recall →", use_container_width=True):
                    st.session_state.ar_phase = "free_recall"
                    st.rerun()
            else:
                st.error(f"Could not read file: {content}")

    # ── PHASE 2: Free Recall ──
    elif st.session_state.ar_phase == "free_recall":
        if st.session_state.ar_round == 1:
            st.markdown("""<div class="recall-phase-box"><div class="recall-phase-title">✍️ Step 2: Write Everything You Remember</div>
                <div class="ebl-phase-desc">Without looking at your notes, write down everything you can remember from your study material.
                Don't worry about order or completeness — just get everything out of your head. This is the most important step: retrieval is where learning happens.</div></div>""", unsafe_allow_html=True)
        else:
            st.markdown(f"""<div class="recall-phase-box"><div class="recall-phase-title">✍️ Round {st.session_state.ar_round}: Write What You Remember Now</div>
                <div class="ebl-phase-desc">You've reviewed your gaps from the previous round. Now, without looking, write everything you remember — especially the concepts you missed or misunderstood last time.</div></div>""", unsafe_allow_html=True)

        st.markdown(f"""<div class="limitation-notice"><strong>📄 Studying:</strong> {st.session_state.ar_file_name}</div>""", unsafe_allow_html=True)

        free_recall = st.text_area(
            "Write everything you remember (no peeking!):",
            height=300,
            placeholder="Start writing everything you can recall from the material...\n\nDon't worry about perfect wording — just get the concepts, facts, mechanisms, and relationships down.",
            key=f"recall_input_round_{st.session_state.ar_round}",
        )

        if st.button("Submit my recall →", use_container_width=True, disabled=not free_recall):
            st.session_state.ar_free_recall = free_recall
            st.session_state.ar_analysis = None
            st.session_state.ar_phase = "gap_report_1"
            st.rerun()

    # ── PHASE 3: Gap Report 1 (after free recall) ──
    elif st.session_state.ar_phase == "gap_report_1":
        st.markdown("""<div class="recall-phase-box"><div class="recall-phase-title">📊 Step 3: What You Missed & Misunderstood</div>
            <div class="ebl-phase-desc">The AI has compared your recall against the original material. Here's what you got right, what you got wrong, and what you missed completely.</div></div>""", unsafe_allow_html=True)

        # Run analysis if not done yet
        if st.session_state.ar_analysis is None:
            with st.spinner("Analysing your recall against the study material..."):
                analysis_messages = [{"role": "user", "content": f"""## ORIGINAL STUDY MATERIAL:
{st.session_state.ar_study_material[:8000]}

## STUDENT'S FREE RECALL (Round {st.session_state.ar_round}):
{st.session_state.ar_free_recall}

Analyse this recall attempt against the original material. Respond ONLY with valid JSON."""}]
                result = call_claude_json(analysis_messages, ACTIVE_RECALL_ANALYSIS_PROMPT)
                if result:
                    st.session_state.ar_analysis = result
                    st.session_state.ar_history.append({"round": st.session_state.ar_round, "phase": "recall", "score": result.get("round_score", 0)})
                    st.rerun()
                else:
                    st.error("Analysis failed. Please try again.")
                    if st.button("Retry Analysis"):
                        st.rerun()

        if st.session_state.ar_analysis:
            _render_gap_report(st.session_state.ar_analysis, "Based on your free recall")

            st.markdown("<br>", unsafe_allow_html=True)

            score = st.session_state.ar_analysis.get("round_score", 0)
            if score >= 95:
                st.markdown("""<div class="reflection-box"><h4>🎉 Outstanding!</h4>
                    <p>You've demonstrated comprehensive recall of this material. Consider discussing the trickiest concepts with your tutor to deepen your understanding even further.</p></div>""", unsafe_allow_html=True)
            else:
                st.markdown("""<div class="reflection-box"><h4>📋 What's next?</h4>
                    <p>Now you can see your gaps. Let's test you with targeted questions focusing on what you missed and misunderstood — then we'll show you an updated gap report.</p></div>""", unsafe_allow_html=True)
                if st.button("❓ Test me with questions →", use_container_width=True):
                    st.session_state.ar_questions = None
                    st.session_state.ar_answers = {}
                    st.session_state.ar_phase = "questions"
                    st.rerun()

    # ── PHASE 4: Questions ──
    elif st.session_state.ar_phase == "questions":
        st.markdown(f"""<div class="recall-phase-box"><div class="recall-phase-title">❓ Step 4: Targeted Questions</div>
            <div class="ebl-phase-desc">These questions focus on your weakest areas first — the concepts you missed or misunderstood in your recall. If you don't know an answer, select "I don't know" — honest gaps are more useful than guesses.</div></div>""", unsafe_allow_html=True)

        # Generate questions if not yet generated
        if st.session_state.ar_questions is None:
            with st.spinner("Generating targeted questions based on your gaps..."):
                q_messages = [{"role": "user", "content": f"""## ORIGINAL STUDY MATERIAL:
{st.session_state.ar_study_material[:6000]}

## KNOWLEDGE GAP ANALYSIS:
{json.dumps(st.session_state.ar_analysis, indent=2)}

## ROUND: {st.session_state.ar_round}

Generate targeted questions prioritising missed and misunderstood concepts. Respond ONLY with valid JSON."""}]
                result = call_claude_json(q_messages, ACTIVE_RECALL_QUESTIONS_PROMPT)
                if result:
                    st.session_state.ar_questions = result
                    st.rerun()
                else:
                    st.error("Question generation failed.")
                    if st.button("Retry"):
                        st.rerun()

        if st.session_state.ar_questions:
            qs = st.session_state.ar_questions
            st.markdown(f"""<div class="reflection-box"><h4>📋 Focus for this round</h4><p>{qs.get('focus_message', 'Answer honestly. Select "I don\'t know" if unsure.')}</p></div>""", unsafe_allow_html=True)

            questions = qs.get("questions", [])

            with st.form("questions_form"):
                for q in questions:
                    qid = str(q.get("id", ""))
                    gap = q.get("gap_type", "")
                    gap_icon = {"missed": "❌", "misunderstood": "⚠️", "understood": "✅"}.get(gap, "")

                    st.markdown(f"**{gap_icon} Q{qid}: {q['question']}**")
                    st.markdown(f"<span style='font-size:0.72rem;color:var(--text-muted);'>Testing: {q.get('concept', '')} · {q.get('difficulty', '')}</span>", unsafe_allow_html=True)

                    answer = st.text_area(
                        "Your answer:",
                        key=f"q_answer_{st.session_state.ar_round}_{qid}",
                        height=80,
                        placeholder="Type your answer here, or leave blank and tick 'I don't know' below",
                    )
                    idk = st.checkbox("🚫 I don't know", key=f"q_idk_{st.session_state.ar_round}_{qid}")

                    st.session_state.ar_answers[qid] = "I DON'T KNOW" if idk else answer
                    st.markdown("---")

                submitted = st.form_submit_button("Submit all answers →", use_container_width=True)

            if submitted:
                # Build answers text
                answers_text = ""
                for q in questions:
                    qid = str(q.get("id", ""))
                    ans = st.session_state.ar_answers.get(qid, "")
                    answers_text += f"Q: {q['question']}\nA: {ans}\n\n"

                st.session_state.ar_question_answers = answers_text
                st.session_state.ar_graded_answers = None
                st.session_state.ar_phase = "q_feedback"
                st.rerun()

    # ── PHASE 4b: Question Feedback ──
    elif st.session_state.ar_phase == "q_feedback":
        st.markdown(f"""<div class="recall-phase-box"><div class="recall-phase-title">📝 Step 4b: Question Feedback</div>
            <div class="ebl-phase-desc">Here's how you did on each question — what was right, what was wrong, and the correct answers so you can learn from each one.</div></div>""", unsafe_allow_html=True)

        # Grade answers if not done yet
        if not st.session_state.get("ar_graded_answers"):
            with st.spinner("Grading your answers..."):
                # Build question list for grading
                questions_for_grading = ""
                qs_data = st.session_state.ar_questions.get("questions", [])
                for q in qs_data:
                    qid = str(q.get("id", ""))
                    ans = st.session_state.ar_answers.get(qid, "")
                    questions_for_grading += f"Question {qid}: {q['question']}\nConcept: {q.get('concept', '')}\nStudent Answer: {ans}\n\n"

                grade_messages = [{"role": "user", "content": f"""## ORIGINAL STUDY MATERIAL:
{st.session_state.ar_study_material[:6000]}

## INITIAL GAP ANALYSIS (from free recall):
Understood: {json.dumps(st.session_state.ar_analysis.get('understood', []), indent=2)}
Misunderstood: {json.dumps(st.session_state.ar_analysis.get('misunderstood', []), indent=2)}
Missed: {json.dumps(st.session_state.ar_analysis.get('missed', []), indent=2)}

## QUESTIONS AND STUDENT ANSWERS:
{questions_for_grading}

Grade each answer against the study material. Respond ONLY with valid JSON."""}]

                result = call_claude_json(grade_messages, ACTIVE_RECALL_GRADE_ANSWERS_PROMPT)
                if result:
                    st.session_state.ar_graded_answers = result
                    st.rerun()
                else:
                    st.error("Grading failed. Please try again.")
                    if st.button("Retry Grading"):
                        st.rerun()

        if st.session_state.get("ar_graded_answers"):
            graded = st.session_state.ar_graded_answers

            # Summary stats
            correct = graded.get("questions_correct", 0)
            partial = graded.get("questions_partial", 0)
            incorrect = graded.get("questions_incorrect", 0)
            idk = graded.get("questions_idk", 0)
            total_q = graded.get("total_questions", 0)

            st.markdown(f"""<div style="text-align:center;margin:1rem 0;">
                <span style="font-size:1.4rem;margin:0 0.8rem;">✅ {correct}</span>
                <span style="font-size:1.4rem;margin:0 0.8rem;">⚠️ {partial}</span>
                <span style="font-size:1.4rem;margin:0 0.8rem;">❌ {incorrect}</span>
                <span style="font-size:1.4rem;margin:0 0.8rem;">🚫 {idk}</span>
                <div style="color:var(--text-muted);font-size:0.82rem;margin-top:0.4rem;">out of {total_q} questions</div>
            </div>""", unsafe_allow_html=True)

            st.markdown(f"""<div style="color:var(--text-secondary);text-align:center;margin-bottom:1.5rem;font-size:0.88rem;">
                {graded.get('overall_feedback', '')}</div>""", unsafe_allow_html=True)

            # Show each graded answer
            for ga in graded.get("graded_answers", []):
                verdict = ga.get("verdict", "")
                if verdict == "correct":
                    icon, border_color, label = "✅", "#5FCC8F", "CORRECT"
                elif verdict == "partially_correct":
                    icon, border_color, label = "⚠️", "#E8C067", "PARTIALLY CORRECT"
                elif verdict == "i_dont_know":
                    icon, border_color, label = "🚫", "#99B3A5", "I DON'T KNOW"
                else:
                    icon, border_color, label = "❌", "#E06060", "INCORRECT"

                improvement = ga.get("improvement_status", "")
                imp_badge = ""
                if improvement == "improved":
                    imp_badge = '<span style="font-size:0.7rem;color:#5FCC8F;font-weight:600;margin-left:8px;">↑ IMPROVED</span>'
                elif improvement == "regression":
                    imp_badge = '<span style="font-size:0.7rem;color:#E06060;font-weight:600;margin-left:8px;">↓ REGRESSION</span>'

                st.markdown(f"""<div style="background:var(--bg-card);border:1px solid var(--border);border-left:4px solid {border_color};
                    border-radius:10px;padding:1rem 1.2rem;margin-bottom:1rem;">
                    <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:0.5rem;">
                        <strong style="color:var(--text-primary);">{icon} Q{ga.get('question_id','')}: {ga.get('question','')}</strong>
                        <span style="font-size:0.72rem;font-weight:700;color:{border_color};">{label}{imp_badge}</span>
                    </div>
                    <div style="font-size:0.85rem;color:var(--text-secondary);margin-bottom:0.4rem;">
                        <strong>Your answer:</strong> {ga.get('student_answer', 'No answer') or 'No answer'}
                    </div>
                    <div style="font-size:0.85rem;color:var(--success);margin-bottom:0.4rem;">
                        <strong>Correct answer:</strong> {ga.get('correct_answer', '')}
                    </div>
                    <div style="font-size:0.82rem;color:var(--text-muted);">
                        {ga.get('explanation', '')}
                    </div>
                </div>""", unsafe_allow_html=True)

            st.markdown("<br>", unsafe_allow_html=True)

            st.markdown("""<div class="reflection-box"><h4>📊 What's next?</h4>
                <p>Now let's combine everything — your initial recall AND your question answers — to give you a complete updated picture of where you stand.</p></div>""", unsafe_allow_html=True)

            if st.button("📊 Show updated gap report →", use_container_width=True):
                st.session_state.ar_analysis_post_questions = None
                st.session_state.ar_phase = "gap_report_2"
                st.rerun()

    # ── PHASE 5: Gap Report 2 (after questions) ──
    elif st.session_state.ar_phase == "gap_report_2":
        st.markdown(f"""<div class="recall-phase-box"><div class="recall-phase-title">📊 Step 5: Updated Gap Report</div>
            <div class="ebl-phase-desc">This combines EVERYTHING — your initial free recall AND your question answers — to give you the most accurate picture of what you know and what still needs work.</div></div>""", unsafe_allow_html=True)

        # Run updated analysis combining original recall + question answers
        if not st.session_state.get("ar_analysis_post_questions"):
            with st.spinner("Building your complete knowledge picture..."):
                # Build a comprehensive summary of all student knowledge demonstrated
                # Include the ORIGINAL free recall (not modified) + graded question results
                graded = st.session_state.get("ar_graded_answers", {})
                graded_summary = ""
                for ga in graded.get("graded_answers", []):
                    v = ga.get("verdict", "")
                    graded_summary += f"Q: {ga.get('question', '')} | Concept: {ga.get('concept', '')} | "
                    graded_summary += f"Student said: {ga.get('student_answer', '')} | Verdict: {v}\n"

                analysis_messages = [{"role": "user", "content": f"""## ORIGINAL STUDY MATERIAL:
{st.session_state.ar_study_material[:8000]}

## PART 1 — STUDENT'S INITIAL FREE RECALL (what they wrote from memory BEFORE any questions):
{st.session_state.ar_free_recall}

## PART 2 — INITIAL GAP ANALYSIS (what the AI found from the free recall):
Total concepts: {st.session_state.ar_analysis.get('total_concepts', 0)}
Understood: {st.session_state.ar_analysis.get('understood_count', 0)}
Misunderstood: {st.session_state.ar_analysis.get('misunderstood_count', 0)}
Missed: {st.session_state.ar_analysis.get('missed_count', 0)}

Understood concepts: {json.dumps([c['concept'] for c in st.session_state.ar_analysis.get('understood', [])], indent=2)}
Misunderstood concepts: {json.dumps([c['concept'] for c in st.session_state.ar_analysis.get('misunderstood', [])], indent=2)}
Missed concepts: {json.dumps([c['concept'] for c in st.session_state.ar_analysis.get('missed', [])], indent=2)}

## PART 3 — GRADED QUESTION ANSWERS (how the student did on targeted questions AFTER seeing their gaps):
{graded_summary}

## YOUR TASK:
Create a COMBINED analysis that accounts for BOTH the free recall AND the question answers together.
- If a concept was UNDERSTOOD in free recall, it stays understood (unless they got it wrong in questions — then it's misunderstood)
- If a concept was MISSED in free recall but answered CORRECTLY in questions, it moves to understood
- If a concept was MISUNDERSTOOD in free recall but answered CORRECTLY in questions, it moves to understood
- If a concept was MISSED in free recall and answered INCORRECTLY or "I DON'T KNOW" in questions, it stays missed
- If a concept was MISUNDERSTOOD in free recall and answered INCORRECTLY in questions, it stays misunderstood
- Concepts not tested by questions keep their original status from the free recall

The round_score should reflect the student's TOTAL knowledge demonstrated across both free recall and questions combined.

Respond ONLY with valid JSON."""}]
                result = call_claude_json(analysis_messages, ACTIVE_RECALL_ANALYSIS_PROMPT)
                if result:
                    st.session_state.ar_analysis_post_questions = result
                    st.session_state.ar_history.append({"round": st.session_state.ar_round, "phase": "questions", "score": result.get("round_score", 0)})
                    st.rerun()
                else:
                    st.error("Analysis failed. Please try again.")
                    if st.button("Retry Analysis"):
                        st.rerun()

        if st.session_state.get("ar_analysis_post_questions"):
            post_q = st.session_state.ar_analysis_post_questions
            pre_q = st.session_state.ar_analysis

            # Show improvement comparison
            pre_score = pre_q.get("round_score", 0)
            post_score = post_q.get("round_score", 0)
            diff = post_score - pre_score

            if diff > 0:
                st.markdown(f"""<div class="reflection-box"><h4>📈 Progress within this round</h4>
                    <p>Your score improved from <strong>{pre_score}%</strong> (free recall alone) to <strong>{post_score}%</strong> (after questions) — that's a <strong>+{diff}%</strong> improvement. The questions helped you demonstrate knowledge you couldn't access freely.</p></div>""", unsafe_allow_html=True)
            elif diff == 0:
                st.markdown(f"""<div class="reflection-box"><h4>📊 Consistent performance</h4>
                    <p>Your score held steady at <strong>{post_score}%</strong>. The questions confirmed what your free recall showed.</p></div>""", unsafe_allow_html=True)
            else:
                st.markdown(f"""<div class="reflection-box"><h4>📊 Results</h4>
                    <p>Your combined score is <strong>{post_score}%</strong>. Some question answers revealed gaps that weren't apparent in your free recall — that's useful to know.</p></div>""", unsafe_allow_html=True)

            _render_gap_report(post_q, "Combined: free recall + question answers")

            # Full round history
            if len(st.session_state.ar_history) > 1:
                st.markdown("**📈 Progress Across All Rounds:**")
                for h in st.session_state.ar_history:
                    bar_w = h['score']
                    bar_c = "#5FCC8F" if h['score'] >= 80 else "#E8C067" if h['score'] >= 60 else "#E06060"
                    phase_label = "Recall" if h.get('phase') == 'recall' else "After Qs"
                    st.markdown(f"""<div style="display:flex;align-items:center;margin:0.3rem 0;font-size:0.82rem;">
                        <span style="width:100px;color:var(--text-muted);">R{h['round']} · {phase_label}</span>
                        <div style="flex:1;height:8px;background:rgba(255,255,255,0.06);border-radius:4px;overflow:hidden;">
                        <div style="width:{bar_w}%;height:100%;background:{bar_c};border-radius:4px;"></div></div>
                        <span style="width:40px;text-align:right;color:var(--text-secondary);font-weight:600;">{h['score']}%</span></div>""", unsafe_allow_html=True)

            st.markdown("<br>", unsafe_allow_html=True)

            if post_score >= 95:
                st.markdown("""<div class="reflection-box"><h4>🎉 Outstanding!</h4>
                    <p>You've demonstrated comprehensive recall of this material. Consider discussing the trickiest concepts with your tutor to deepen your understanding even further.</p></div>""", unsafe_allow_html=True)
            else:
                st.markdown("""<div class="reflection-box"><h4>🔄 Ready for another round?</h4>
                    <p>Review the gaps above — especially the ❌ missed and ⚠️ misunderstood sections. When you feel ready, start a new round. Each cycle strengthens your recall pathways.</p></div>""", unsafe_allow_html=True)

                if st.button("🔄 Start next round →", use_container_width=True):
                    st.session_state.ar_round += 1
                    st.session_state.ar_phase = "free_recall"
                    st.session_state.ar_free_recall = None
                    st.session_state.ar_analysis = None
                    st.session_state.ar_analysis_post_questions = None
                    st.session_state.ar_questions = None
                    st.session_state.ar_answers = {}
                    st.session_state.ar_question_answers = None
                    st.session_state.ar_graded_answers = None
                    st.rerun()


# ─── Feedback Mode ───
elif st.session_state.mode == "feedback":

    st.markdown(f"""<div style="margin-bottom:1.5rem;">
        <span style="font-family:'DM Serif Display',serif;font-size:1.4rem;color:var(--text-primary);">
        💬 Platform Feedback</span></div>""", unsafe_allow_html=True)

    st.markdown("""<div class="ebl-phase">
        <div class="ebl-phase-title">Help Us Improve DentEdTech™</div>
        <div class="ebl-phase-desc">Your feedback is essential for improving this platform's performance and credibility.
        Please be as honest and specific as possible — constructive criticism is especially valuable.
        Your responses will be sent directly to the development team.</div>
    </div>""", unsafe_allow_html=True)

    if st.session_state.feedback_submitted:
        st.markdown("""<div class="reflection-box">
            <h4>✅ Thank you for your feedback!</h4>
            <p>Your feedback has been sent successfully to the DentEdTech™ development team.
            Your input directly shapes future improvements to this platform.</p>
        </div>""", unsafe_allow_html=True)
        if st.button("← Back to Mode Selection", use_container_width=True, key="btn_feedback_back"):
            st.session_state.mode = None
            st.session_state.feedback_submitted = False
            st.rerun()
    else:
        st.markdown("##### About You *(optional)*")
        fb_name = st.text_input("Your name", placeholder="Optional — you can remain anonymous", key="fb_name")
        fb_email = st.text_input("Your email", placeholder="Optional — only if you'd like a response", key="fb_email")
        fb_discipline = st.selectbox("Your discipline", ["Dentistry", "Medicine", "Pharmacology", "Other"], index=0, key="fb_disc")
        fb_year = st.selectbox("Year of study", ["Year 1", "Year 2", "Year 3", "Year 4", "Year 5", "Postgraduate", "Faculty/Staff"], index=2, key="fb_year")

        st.markdown("---")
        st.markdown("##### Which features did you use? *(select all that apply)*")
        used_evidence = st.checkbox("🔬 Evidence-Based Knowledge", key="fb_used_evidence")
        used_ebl = st.checkbox("🧭 Enquiry-Based Learning", key="fb_used_ebl")
        used_video = st.checkbox("🎥 Clinical Video Trust Engine", key="fb_used_video")
        used_recall = st.checkbox("🧠 Active Recall", key="fb_used_recall")

        st.markdown("---")
        st.markdown("##### Your Feedback")

        fb_worked = st.text_area("What worked well?", height=120,
            placeholder="Which features were most useful? What did you enjoy? What helped your learning?", key="fb_worked")
        fb_improve = st.text_area("What needs improvement?", height=120,
            placeholder="What was confusing, frustrating, or unhelpful? What didn't work as expected?", key="fb_improve")
        fb_suggestions = st.text_area("How could we improve it? Any suggestions or new features you'd like?", height=120,
            placeholder="Be as specific as possible — your ideas directly shape the next version.", key="fb_suggestions")

        st.markdown("---")
        st.markdown("##### Overall Experience")

        fb_rating = st.select_slider("Overall rating",
            options=["1 — Poor", "2 — Below Average", "3 — Average", "4 — Good", "5 — Excellent"],
            value="3 — Average", key="fb_rating")
        fb_recommend = st.radio("Would you recommend DentEdTech™ to a fellow student?",
            options=["Definitely yes", "Probably yes", "Not sure", "Probably not", "Definitely not"],
            index=2, key="fb_recommend")
        fb_additional = st.text_area("Anything else you'd like to tell us?", height=80,
            placeholder="Any additional comments, concerns, or praise.", key="fb_additional")

        if st.button("Submit Feedback →", use_container_width=True, key="btn_submit_feedback"):
            # Build features list
            features_used = []
            if used_evidence: features_used.append("Evidence-Based Knowledge")
            if used_ebl: features_used.append("Enquiry-Based Learning")
            if used_video: features_used.append("Clinical Video Trust Engine")
            if used_recall: features_used.append("Active Recall")

            # Build the full message
            message_body = f"""DentEdTech™ Platform Feedback
================================
Submitted: {datetime.now().strftime('%Y-%m-%d %H:%M')}

ABOUT THE USER
Name: {fb_name or 'Anonymous'}
Email: {fb_email or 'Not provided'}
Discipline: {fb_discipline}
Year: {fb_year}

FEATURES USED
{', '.join(features_used) if features_used else 'None selected'}

WHAT WORKED WELL
{fb_worked or 'No response'}

WHAT NEEDS IMPROVEMENT
{fb_improve or 'No response'}

SUGGESTIONS FOR IMPROVEMENT
{fb_suggestions or 'No response'}

OVERALL RATING
{fb_rating}

WOULD RECOMMEND
{fb_recommend}

ADDITIONAL COMMENTS
{fb_additional or 'No response'}

================================
Sent from DentEdTech™ Evidence Engine"""

            # Send via Gmail SMTP (direct, no third-party service)
            try:
                gmail_user = st.secrets["GMAIL_ADDRESS"]
                gmail_app_password = st.secrets["GMAIL_APP_PASSWORD"]
            except (KeyError, FileNotFoundError):
                gmail_user = None
                gmail_app_password = None

            if not gmail_user or not gmail_app_password:
                st.error("Feedback email not configured. Please add GMAIL_ADDRESS and GMAIL_APP_PASSWORD to Streamlit secrets.")
            else:
                recipient = "ayman.khalifah@manchester.ac.uk"
                subject = f"DentEdTech™ Feedback — {fb_discipline} {fb_year} — {datetime.now().strftime('%d/%m/%Y')}"

                msg = MIMEMultipart()
                msg["From"] = gmail_user
                msg["To"] = recipient
                msg["Subject"] = subject
                msg["Reply-To"] = fb_email or gmail_user
                msg.attach(MIMEText(message_body, "plain"))

                try:
                    with smtplib.SMTP("smtp.gmail.com", 587, timeout=15) as server:
                        server.starttls()
                        server.login(gmail_user, gmail_app_password)
                        server.sendmail(gmail_user, recipient, msg.as_string())
                    st.session_state.feedback_submitted = True
                    st.rerun()
                except smtplib.SMTPAuthenticationError:
                    st.error("Gmail authentication failed. Please check your GMAIL_ADDRESS and GMAIL_APP_PASSWORD in Streamlit secrets.")
                except Exception as e:
                    st.error(f"Could not send feedback: {str(e)}")
                    st.markdown(f"""<div class="limitation-notice">
                        <strong>Alternative:</strong> Please copy the feedback below and email it manually to
                        <strong>ayman.khalifah@manchester.ac.uk</strong>
                    </div>""", unsafe_allow_html=True)
                    st.code(message_body, language=None)
